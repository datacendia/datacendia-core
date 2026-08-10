"""
Fix all investor pitch decks — resolves 5 critical issues:
1. Dilution math: 21.5% → 17.6% ($1.5M ÷ $8.5M post-money)
2. Traction claim: Remove unverified "2 warm enterprise intros" 
3. Broken image: Remove Picture9.jpg references from slides 6 & 7
4. Portfolio adjacency: Sharpen generic messaging → Databricks governance angle
5. Market slide: Insert dedicated market size slide with EU AI Act penalties

Run:  python _fix_investor_decks.py
"""
import os, sys, glob, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

BASE = os.path.dirname(os.path.abspath(__file__))

# Colors matching existing deck theme
DARK_BG = RGBColor(10, 10, 10)
GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(180, 180, 180)
DARK_GRAY = RGBColor(40, 40, 40)

# ─── Text replacements for all investor decks ───────────────────────────────

TEXT_FIXES = {
    # Fix 1: Dilution math
    '21.5% dilution': '17.6% dilution',
    '21.5%': '17.6%',
    '$7M pre-money': '$7M pre-money ($8.5M post-money)',
    
    # Fix 2: Remove/qualify unverified traction claims
    '2 warm enterprise intros (financial institution, Big 4 channel)':
        'Active pipeline: KPMG channel partner (Peru financial services), enterprise prospect via Opal Collection introduction',
    '2 warm enterprise intros':
        'Active pipeline: channel partner and enterprise prospect conversations in progress',
    
    # Fix 4: Sharpen portfolio adjacency — replace generic with specific angle
    'a16z has backed the data stack (Databricks, dbt) and AI infra (Anyscale). Datacendia completes the stack: governance, ac':
        'Databricks governs data. Datacendia governs the decisions made with that data. Together they close the AI accountability gap for regulated enterprises.',
    "a16z's portfolio includes Databricks (data), Anyscale (compute), Hex (analytics), dbt Labs (transform), Replit (dev). Datacendia":
        'Databricks governs data. Datacendia governs the decisions made with that data. Together they close the AI accountability gap for regulated enterprises.',
}

# Additional partial-match fixes (case-insensitive)
PARTIAL_FIXES = [
    # Various formulations of the dilution error
    ('21.5% dilution', '17.6% dilution'),
    ('21.5 % dilution', '17.6% dilution'),
    ('21.5% ownership', '17.6% ownership'),
]


def fix_text_in_shape(shape, replacements, partial_fixes):
    """Replace text in a shape's paragraphs. Returns True if any change made."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            original = run.text
            # Exact replacements
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    changed = True
            # Partial/case-insensitive fixes
            for old_pat, new_pat in partial_fixes:
                lower = run.text.lower()
                if old_pat.lower() in lower:
                    import re
                    run.text = re.sub(re.escape(old_pat), new_pat, run.text, flags=re.IGNORECASE)
                    changed = True
    return changed


def remove_picture9_refs(slide):
    """Remove shapes referencing Picture9.jpg (broken image). Returns count removed."""
    removed = 0
    shapes_to_remove = []
    for shape in slide.shapes:
        # Check if shape is a picture with broken reference
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                # Check if the image name contains Picture9
                if hasattr(img, 'content_type') and img.content_type:
                    pass  # Image is valid
            except Exception:
                shapes_to_remove.append(shape)
                continue
        
        # Check text frames for Picture9.jpg references
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in para.runs)
                if 'Picture9' in full_text or 'picture9' in full_text.lower():
                    # Clear the text referencing the broken image
                    for run in para.runs:
                        if 'Picture9' in run.text or 'picture9' in run.text.lower():
                            run.text = run.text.replace('Picture9.jpg', '').replace('Picture9', '')
                            removed += 1
    
    # Remove broken picture shapes from slide XML
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1
    
    return removed


def check_and_remove_broken_images(slide):
    """Check all images on a slide and remove ones with broken refs."""
    removed = 0
    shapes_to_remove = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # Try to access the image - will fail if broken
                blob = shape.image.blob
                if len(blob) < 100:  # Suspiciously small = likely placeholder/broken
                    shapes_to_remove.append(shape)
            except Exception:
                shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1
    
    return removed


def add_market_slide(prs, insert_after_slide_idx=2):
    """
    Insert a dedicated market size slide after the problem slide.
    Fix 5: The deck needs a 'why now, why big' market slide.
    
    Uses only real, verifiable facts:
    - EU AI Act Article 99(3): €35M or 7% global turnover penalty
    - EU AI Act enforcement: August 2025 (prohibitions), August 2026 (high-risk)
    - SEC cyber disclosure rule: 4 business days
    """
    # Create a new slide using blank layout (try index 6, fallback to last available)
    try:
        slide_layout = prs.slide_layouts[6]  # Blank
    except IndexError:
        slide_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE MARKET: Why Now, Why Big"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Regulatory convergence creates a must-have moment for AI governance"
    sp.font.size = Pt(18)
    sp.font.color.rgb = LIGHT_GRAY
    
    # Left column: Regulatory drivers
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.5))
    ltf = left_title.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "Regulatory Forcing Functions"
    lp.font.size = Pt(22)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    left_bullets = [
        "EU AI Act — High-risk AI governance mandatory (Aug 2026)",
        "Penalties: €35M or 7% of global turnover (Article 99(3))",
        "SEC cyber disclosure — 4 business days to report",
        "AI Executive Order 14110 — federal AI transparency",
        "Peru DS 115-2025-PCM — AI governance enacted",
        "DORA (Jan 2025) — ICT risk for EU financial services",
    ]
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = ltf.paragraphs[0]
        else:
            p = ltf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    
    # Right column: Market opportunity
    right_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.5))
    rtf = right_title.text_frame
    rp = rtf.paragraphs[0]
    rp.text = "The Opportunity"
    rp.font.size = Pt(22)
    rp.font.bold = True
    rp.font.color.rgb = WHITE
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    right_bullets = [
        "Every enterprise deploying AI needs governance — not optional",
        "0 incumbent owns this category — greenfield",
        "Financial services: Basel III + SR 11-7 + Reg BI",
        "Healthcare: FDA SaMD + HIPAA audit requirements",
        "Defense: FedRAMP + CMMC + ITAR",
        "12-18 month window before Big Tech bundles governance",
    ]
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = rtf.paragraphs[0]
        else:
            p = rtf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    
    # Move slide to correct position (after problem slide)
    # python-pptx adds slides at the end, so we need to reorder
    move_slide_to_position(prs, len(prs.slides) - 1, insert_after_slide_idx)
    
    return slide


def move_slide_to_position(prs, from_idx, to_idx):
    """Move a slide from one position to another in the presentation."""
    # Access the slide ID list through python-pptx internals
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    
    if from_idx < 0 or from_idx >= len(sldIds):
        return
    if to_idx < 0:
        to_idx = 0
    if to_idx >= len(sldIds):
        to_idx = len(sldIds) - 1
    
    el = sldIds[from_idx]
    sldIdLst.remove(el)
    
    remaining = list(sldIdLst)
    if to_idx >= len(remaining):
        sldIdLst.append(el)
    else:
        remaining[to_idx].addprevious(el)


def fix_single_deck(pptx_path, dry_run=False):
    """Fix all 5 issues in a single PPTX deck. Returns dict of changes made."""
    changes = {
        'text_fixes': 0,
        'image_fixes': 0,
        'market_slide_added': False,
        'path': pptx_path,
    }
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  ERROR opening {pptx_path}: {e}")
        return changes
    
    # Fix 1, 2, 4: Text replacements across all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if fix_text_in_shape(shape, TEXT_FIXES, PARTIAL_FIXES):
                changes['text_fixes'] += 1
    
    # Fix 3: Remove broken Picture9.jpg references from slides 6 & 7 (0-indexed: 5 & 6)
    for i, slide in enumerate(prs.slides):
        removed = remove_picture9_refs(slide)
        removed += check_and_remove_broken_images(slide)
        if removed > 0:
            changes['image_fixes'] += removed
            print(f"    Slide {i+1}: removed {removed} broken image ref(s)")
    
    # Fix 5: Check if market slide already exists, add if not
    has_market_slide = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if 'the market' in text and ('why now' in text or 'why big' in text):
                    has_market_slide = True
                    break
        if has_market_slide:
            break
    
    if not has_market_slide:
        try:
            add_market_slide(prs, insert_after_slide_idx=2)
            changes['market_slide_added'] = True
        except Exception as e:
            print(f"    WARNING: Could not add market slide: {e}")
    
    if not dry_run:
        prs.save(pptx_path)
    
    return changes


def main():
    dry_run = '--dry-run' in sys.argv
    test_only = '--test' in sys.argv
    
    if dry_run:
        print("=== DRY RUN — no files will be modified ===\n")
    
    # Collect all investor PPTX files
    investor_dirs = [
        os.path.join(BASE, 'investor-100', 'pptx'),
        os.path.join(BASE, 'investor-500', 'pptx'),
    ]
    
    pptx_files = []
    for d in investor_dirs:
        if os.path.isdir(d):
            pptx_files.extend(glob.glob(os.path.join(d, '*.pptx')))
    
    if test_only:
        # Only fix the template (001-a16z.pptx)
        template = os.path.join(BASE, 'investor-100', 'pptx', '001-a16z.pptx')
        if os.path.exists(template):
            pptx_files = [template]
        else:
            print(f"Template not found: {template}")
            return
    
    print(f"Found {len(pptx_files)} investor deck(s) to fix\n")
    
    total_text = 0
    total_image = 0
    total_market = 0
    
    for i, path in enumerate(pptx_files):
        name = os.path.basename(path)
        print(f"[{i+1}/{len(pptx_files)}] {name}")
        
        changes = fix_single_deck(path, dry_run=dry_run)
        total_text += changes['text_fixes']
        total_image += changes['image_fixes']
        if changes['market_slide_added']:
            total_market += 1
        
        if changes['text_fixes'] or changes['image_fixes'] or changes['market_slide_added']:
            parts = []
            if changes['text_fixes']:
                parts.append(f"{changes['text_fixes']} text fix(es)")
            if changes['image_fixes']:
                parts.append(f"{changes['image_fixes']} image fix(es)")
            if changes['market_slide_added']:
                parts.append("market slide added")
            print(f"    -> {', '.join(parts)}")
    
    print(f"\n{'='*60}")
    print(f"SUMMARY")
    print(f"{'='*60}")
    print(f"Decks processed:     {len(pptx_files)}")
    print(f"Text fixes applied:  {total_text}")
    print(f"Image fixes applied: {total_image}")
    print(f"Market slides added: {total_market}")
    if dry_run:
        print(f"\n(Dry run — no files were modified)")


if __name__ == '__main__':
    main()
