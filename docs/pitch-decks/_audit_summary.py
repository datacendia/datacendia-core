"""Summarize: for each deck CATEGORY, show unique image slide content patterns."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os, glob
from collections import defaultdict

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'
folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 
           'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

for folder in folders:
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    
    print(f'\n{"="*80}')
    print(f'{folder}/ --- {len(files)} decks')
    print(f'{"="*80}')
    
    # Collect unique patterns: (slide_position, image_slide_index, content_signature)
    patterns = defaultdict(list)
    
    for f in files:
        try:
            prs = Presentation(f)
            name = os.path.basename(f)
            total = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                has_img = False
                img_pos = None
                texts = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        has_img = True
                        img_pos = (round(shape.left/914400,1), round(shape.top/914400,1),
                                   round(shape.width/914400,1), round(shape.height/914400,1))
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and t != 'DATACENDIA' and 'datacendia.com' not in t.lower() and '2024' not in t and '2026' not in t:
                                texts.append(t)
                
                if has_img:
                    # Key: which image slide is this (1st or 2nd image in deck), what's on it
                    content_key = ' | '.join(texts[:3])  # first 3 meaningful texts
                    key = f'S{i+1}/{total} img@({img_pos[0]},{img_pos[1]}) {img_pos[2]}x{img_pos[3]}in: {content_key[:120]}'
                    patterns[key].append(name)
        except:
            pass
    
    for key, decks in sorted(patterns.items()):
        print(f'\n  {key}')
        print(f'    -> {len(decks)} decks: {", ".join(decks[:3])}{"..." if len(decks) > 3 else ""}')
