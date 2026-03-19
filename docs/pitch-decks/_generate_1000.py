"""
Generate pitch decks by cloning an existing template and swapping personalized text.
Uses 001-a16z.pptx as the template for investor decks.
Uses 01-fortune-500.pptx as the template for sales decks.
Uses 01-healthcare.pptx as the template for design-partner decks.

Usage:
  python _generate_1000.py --test       # Generate 1 test deck
  python _generate_1000.py --category investor-500  # Generate one category
  python _generate_1000.py              # Generate all 760 new decks
"""
import copy, os, sys, glob, shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'
SCREENSHOTS = os.path.join(BASE, 'screenshots')

# ─── Template paths ──────────────────────────────────────────────────────────
TEMPLATES = {
    'investor': os.path.join(BASE, 'investor-100', 'pptx', '001-a16z.pptx'),
    'sales': os.path.join(BASE, 'sales', 'pptx', '01-fortune-500.pptx'),
    'design-partner': os.path.join(BASE, 'design-partner', 'pptx', '01-healthcare.pptx'),
    'investor-50': os.path.join(BASE, 'investor-50', 'pptx', '01-banking.pptx'),
    'gamma': os.path.join(BASE, 'gamma', 'pptx', '01-investor-overview.pptx'),
}

# ─── Text replacement mapping ────────────────────────────────────────────────
# For investor-100 template (a16z), the personalized text fields are:

def replace_text_in_shape(shape, replacements):
    """Replace text in a shape's paragraphs using a dict of old->new mappings."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            original = run.text
            for old_text, new_text in replacements.items():
                if old_text.lower() in run.text.lower():
                    # Case-insensitive replace but preserve the run's text case pattern
                    import re
                    run.text = re.sub(re.escape(old_text), new_text, run.text, flags=re.IGNORECASE)
                    changed = True
    return changed


def generate_investor_deck(firm, output_dir):
    """Generate a personalized investor deck from the a16z template."""
    template_path = TEMPLATES['investor']
    prs = Presentation(template_path)
    
    # Replacements: template text -> new text
    name = firm['name']
    name_upper = name.upper()
    slug = firm['slug']
    
    replacements = {
        'ANDREESSEN HOROWITZ (A16Z)': name_upper,
        'Andreessen Horowitz (a16z)': name,
        'Andreessen Horowitz': name,
        'a16z': firm.get('slug_display', slug),
    }
    
    # Firm-specific fields
    old_type_line = 'Venture Capital · Menlo Park, CA · Seed to Growth · $35B+ AUM'
    new_type_line = f'{firm["type"]} · {firm["loc"]} · {firm["aum"]} AUM'
    replacements[old_type_line] = new_type_line
    
    # Thesis
    old_thesis = 'AI governance is the next infrastructure default'
    new_thesis = firm.get('thesis', 'AI governance infrastructure for regulated enterprises')
    replacements[old_thesis] = new_thesis
    
    # Portfolio companies
    old_portfolio = ['Databricks', 'Anyscale', 'Hex', 'dbt Labs', 'Replit']
    new_portfolio = firm.get('portfolio', [])
    
    # AI thesis alignment
    old_alignment = "AI is eating software. a16z's dedicated AI fund backs infrastructure that becomes the default layer for enterprise AI ad"
    new_alignment = f"{name} invests in {firm['type'].lower()} opportunities. Datacendia is the missing governance layer for enterprise AI."
    replacements[old_alignment] = new_alignment[:120]
    
    old_why = "Your Thesis → Our Mission"
    replacements[old_why] = "Your Thesis \u2192 Our Mission"
    
    old_mission = "Datacendia is the missing governance layer for enterprise AI — the infrastructure that sits between AI deployment and re"
    new_mission = "Datacendia is the missing governance layer for enterprise AI \u2014 the infrastructure that sits between AI deployment and regulatory compliance."
    replacements[old_mission] = new_mission[:120]
    
    old_adjacency_text = "a16z has backed the data stack (Databricks, dbt) and AI infra (Anyscale). Datacendia completes the stack: governance, ac"
    if new_portfolio:
        portfolio_str = ', '.join(new_portfolio[:4])
        new_adjacency = f"{name}'s portfolio includes {portfolio_str}. Datacendia completes the stack with AI governance and audit."
    else:
        new_adjacency = f"Datacendia adds AI governance and cryptographic audit trails to {name}'s portfolio."
    replacements[old_adjacency_text] = new_adjacency[:120]
    
    # Replace portfolio company names in slide 2
    for old_co, new_co in zip(old_portfolio, new_portfolio[:5] if new_portfolio else ['']*5):
        if old_co and new_co:
            replacements[old_co] = new_co
    
    # "Key Angle for" line  
    old_angle = 'Key Angle for Andreessen Horowitz (a16z)'
    new_angle = f'Key Angle for {name}'
    replacements[old_angle] = new_angle
    
    # Sector line
    old_sector = 'Sector: AI Infrastructure'
    new_sector = f'Sector: {firm.get("sector", "AI Infrastructure")}'
    replacements[old_sector] = new_sector
    
    # Next steps personalization
    old_next = 'Personalized for Andreessen Horowitz (a16z) \u2014 AI governance is the next infrastructure default'
    new_next = f'Personalized for {name} \u2014 {new_thesis}'
    replacements[old_next] = new_next[:120]
    
    # Apply replacements to all shapes on all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    # Save
    os.makedirs(output_dir, exist_ok=True)
    num = firm['num']
    filename = f'{str(num).zfill(3)}-{slug}.pptx'
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    return output_path


def generate_sales_deck(target, output_dir):
    """Generate a sales deck from the fortune-500 template."""
    template_path = TEMPLATES['sales']
    prs = Presentation(template_path)
    
    name = target['name']
    replacements = {
        'Fortune 500': name,
        'FORTUNE 500': name.upper(),
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    os.makedirs(output_dir, exist_ok=True)
    num = target['num']
    filename = f'{str(num).zfill(2)}-{target["slug"]}.pptx'
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    return output_path


def generate_design_partner_deck(target, output_dir):
    """Generate a design partner deck from the healthcare template."""
    template_path = TEMPLATES['design-partner']
    prs = Presentation(template_path)
    
    name = target['name']
    replacements = {
        'Healthcare': name,
        'HEALTHCARE': name.upper(),
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    os.makedirs(output_dir, exist_ok=True)
    num = target['num']
    filename = f'{str(num).zfill(2)}-{target["slug"]}.pptx'
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    return output_path


def generate_thematic_deck(target, output_dir, template_key='investor-50'):
    """Generate a thematic investor/regional/enterprise deck."""
    template_path = TEMPLATES.get(template_key, TEMPLATES['investor-50'])
    prs = Presentation(template_path)
    
    name = target['name']
    replacements = {
        'Banking': name,
        'BANKING': name.upper(),
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
    
    os.makedirs(output_dir, exist_ok=True)
    num = target['num']
    filename = f'{str(num).zfill(2)}-{target["slug"]}.pptx'
    if num >= 100:
        filename = f'{str(num).zfill(3)}-{target["slug"]}.pptx'
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)
    return output_path


def embed_screenshots(pptx_path):
    """Replace image placeholders with correct screenshots (same as embed_screenshots.py)."""
    CONTENT_MAP = {
        'multi-agent decision engine': os.path.join(SCREENSHOTS, '02-council-deliberation.png'),
        'decision audit trail': os.path.join(SCREENSHOTS, '05-decisions.png'),
    }
    
    prs = Presentation(pptx_path)
    changes = 0
    
    for slide in prs.slides:
        existing_images = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not existing_images:
            continue
        
        slide_text = ' '.join(
            para.text.strip().lower()
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        
        for keyword, img_path in CONTENT_MAP.items():
            if keyword in slide_text and os.path.exists(img_path):
                old_img = existing_images[0]
                left, top, width, height = old_img.left, old_img.top, old_img.width, old_img.height
                
                for img_shape in existing_images:
                    sp = img_shape._element
                    sp.getparent().remove(sp)
                
                slide.shapes.add_picture(img_path, left, top, width, height)
                changes += 1
                break
    
    if changes > 0:
        prs.save(pptx_path)
    return changes


# ─── Import firm data ────────────────────────────────────────────────────────
from _firms_database import INVESTOR_500_NEW, SALES_EXTRA, DESIGN_PARTNER_EXTRA


# ─── Additional data for remaining categories ────────────────────────────────

# Investor-500: firms 201-500 (300 more to reach 500 total)
# These are generated from a comprehensive list of real global VC/PE firms
INVESTOR_500_BATCH2 = []
_more_firms = [
    # US VC continued (201-250)
    ('201','a-star','A* (A-Star) Partners','Venture Capital','San Francisco, CA','$500M+','us',['AI startups'],'AI-native fund'),
    ('202','amplify','Amplify Partners','Venture Capital','Menlo Park, CA','$1B+','us',['Hashicorp','Datadog','Fastly'],'Developer infrastructure'),
    ('203','bessemer-india','Bessemer Venture Partners (India)','Venture Capital','Bangalore, India','$4B+','asia',['Swiggy','PhonePe','Myntra'],'India enterprise tech'),
    ('204','bold-capital','Bold Capital Partners','Venture Capital','Mountain View, CA','$500M+','us',['Peter Diamandis network'],'Exponential technology'),
    ('205','bowery','Bowery Capital','Venture Capital','New York, NY','$500M+','us',['Nuvocargo','Threekit'],'B2B software seed stage'),
    ('206','bvp-forge','BVP Forge','Growth Equity','San Francisco, CA','$2B+','us',['Twilio','LinkedIn','Yelp'],'Bessemer growth vehicle'),
    ('207','cota','Cota Capital','Venture Capital','San Francisco, CA','$1.5B+','us',['Illumio','Netskope','Sumo Logic'],'Enterprise security and cloud'),
    ('208','costanoa','Costanoa Ventures','Venture Capital','Palo Alto, CA','$500M+','us',['Clari','Lattice','Cloudflare'],'Enterprise seed specialist'),
    ('209','decibel','Decibel Partners','Venture Capital','Menlo Park, CA','$500M+','us',['Abnormal Security','Persona'],'Enterprise software founders fund'),
    ('210','eniac','Eniac Ventures','Venture Capital','New York / SF','$500M+','us',['Flashfood','Flutterwave'],'Pre-seed mobile and enterprise'),
    ('211','fika','Fika Ventures','Venture Capital','Los Angeles, CA','$400M+','us',['Scopely','PatientPop'],'LA-based enterprise tech'),
    ('212','foundry','Foundry Group','Venture Capital','Boulder, CO','$3B+','us',['Fitbit','SendGrid','MakerBot'],'Brad Feld — enterprise and platforms'),
    ('213','greycroft','Greycroft','Venture Capital','New York, NY','$2B+','us',['Venmo','Huffington Post','Braintree'],'Consumer and enterprise investing'),
    ('214','homebrew','Homebrew','Venture Capital','San Francisco, CA','$500M+','us',['Plaid','Gusto','Cruise'],'Seed-stage enterprise leaders'),
    ('215','lerer','Lerer Hippeau','Venture Capital','New York, NY','$1B+','us',['Allbirds','Casper','Buzzfeed'],'NYC ecosystem enterprise tech'),
    ('216','matrix','Matrix Partners','Venture Capital','San Francisco, CA','$5B+','us',['HubSpot','Zendesk','Oculus'],'Multi-stage enterprise VC'),
    ('217','maveron','Maveron','Venture Capital','Seattle, WA','$1.5B+','us',['Zulily','Potbelly','General Assembly'],'Consumer and enterprise — Starbucks origin'),
    ('218','nfx','NFX','Venture Capital','San Francisco, CA','$1B+','us',['Lyft','Trulia','Paylocity'],'Network effects specialists'),
    ('219','norwest','Norwest Venture Partners','Venture Capital','San Jose, CA','$12B+','us',['Workiva','Udemy','Calm'],'Multi-stage enterprise investing'),
    ('220','osage','Osage Venture Partners','Venture Capital','Wayne, PA','$300M+','us',['Sievert Larsen','Centric Software'],'East Coast enterprise B2B'),
    ('221','polaris','Polaris Partners','Venture Capital','Boston, MA','$5B+','us',['Akamai','CarGurus','Quanterix'],'Healthcare and enterprise tech'),
    ('222','rre','RRE Ventures','Venture Capital','New York, NY','$2B+','us',['Venmo','MongoDB','Yext'],'NYC enterprise infrastructure'),
    ('223','shasta','Shasta Ventures','Venture Capital','Menlo Park, CA','$2B+','us',['Nest','Zuora','Anaplan'],'Enterprise and IoT'),
    ('224','sutter-hill','Sutter Hill Ventures','Venture Capital','Palo Alto, CA','$3B+','us',['Snowflake','Pure Storage','Vlocity'],'Created Snowflake — enterprise data'),
    ('225','upfront','Upfront Ventures','Venture Capital','Los Angeles, CA','$2B+','us',['Ring','TrueCar','Overture'],'LA enterprise leader — Mark Suster'),
    # European continued (226-275)
    ('226','balderton-eu','Balderton Capital','Venture Capital','London, UK','$5B+','eu',['Revolut','GoCardless','Citymapper'],'European enterprise leader'),
    ('227','blisce','Blisce','Growth Equity','Paris / NYC','$1B+','eu',['Spotify','Pinterest','Headspace'],'Consumer and enterprise platform growth'),
    ('228','capnamic','Capnamic Ventures','Venture Capital','Cologne, Germany','$300M+','eu',['German B2B startups'],'German enterprise specialist'),
    ('229','cherry-vc','Cherry Ventures','Venture Capital','Berlin, Germany','$1B+','eu',['Flixbus','Auto1','Forto'],'Berlin seed to Series A enterprise'),
    ('230','daphni','Daphni','Venture Capital','Paris, France','$300M+','eu',['Alan','Swile','Sunday'],'Community-driven French VC'),
    ('231','elaia','Elaia Partners','Venture Capital','Paris, France','$500M+','eu',['Criteo','Volterra','Shift Technology'],'French deep tech and enterprise'),
    ('232','eleven-gs','Eleven Global Strategy','Venture Capital','Riyadh / London','$200M+','middle-east',['MENA tech'],'Saudi-European bridge fund'),
    ('233','eurazeo-growth','Eurazeo Growth','Growth Equity','Paris, France','$5B+','eu',['Doctolib','Contentsquare','PayFit'],'French growth stage leader'),
    ('234','fly-ventures','Fly Ventures','Venture Capital','Berlin, Germany','$200M+','eu',['Luminovo','DeepL'],'AI and automation seed fund'),
    ('235','headline','Headline (ex-e.ventures)','Venture Capital','Berlin / San Francisco','$2B+','eu',['Farfetch','Groupon','Sonos'],'Transatlantic enterprise'),
    ('236','holtzbrinck','Holtzbrinck Ventures','Venture Capital','Munich, Germany','$1B+','eu',['Zalando','FlixMobility','Raisin'],'DACH enterprise technology'),
    ('237','isai','ISAI','Venture Capital','Paris, France','$300M+','eu',['Algolia','MeilliSearch'],'French founder-backed VC'),
    ('238','kfund','K Fund','Venture Capital','Madrid, Spain','$200M+','eu',['Jobandtalent','Wallapop'],'Spanish enterprise tech'),
    ('239','kinnevik','Kinnevik','Growth Equity','Stockholm, Sweden','$10B+','eu',['Zalando','Tele2','Betterment'],'Nordic holding — enterprise and fintech'),
    ('240','local-globe','LocalGlobe','Venture Capital','London, UK','$1B+','eu',['Transferwise','Improbable','Tessian'],'London seed and Series A'),
    ('241','mangrove','Mangrove Capital Partners','Venture Capital','Luxembourg','$2B+','eu',['Skype','Wix','Walkme'],'Benelux enterprise technology'),
    ('242','moonfire','Moonfire Ventures','Venture Capital','London, UK','$300M+','eu',['European AI startups'],'AI-focused European seed fund'),
    ('243','new-wave','New Wave','Venture Capital','London, UK','$500M+','eu',['AI and enterprise startups'],'European enterprise AI'),
    ('244','octopus','Octopus Ventures','Venture Capital','London, UK','$2B+','eu',['Zoopla','Secret Escapes','LoveCrafts'],'UK enterprise and consumer'),
    ('245','picus','Picus Capital','Venture Capital','Munich, Germany','$500M+','eu',['German enterprise startups'],'DACH B2B SaaS'),
    ('246','saga-vc','Saga Ventures','Venture Capital','Helsinki, Finland','$200M+','eu',['Nordic AI startups'],'Nordic AI and enterprise'),
    ('247','seedcamp','Seedcamp','Venture Capital','London, UK','$500M+','eu',['Transferwise','UiPath','Revolut'],'Pan-European seed pioneer'),
    ('248','target-global','Target Global','Venture Capital','Berlin / Tel Aviv','$3B+','eu',['Delivery Hero','Auto1'],'European-Israeli enterprise bridge'),
    ('249','tq-ventures','TQ Ventures','Venture Capital','Los Angeles, CA','$500M+','us',['Plaid','Chime','Dutchie'],'Music industry founders — enterprise tech'),
    ('250','ventech','Ventech','Venture Capital','Paris, France','$500M+','eu',['Dailymotion','SigFox','Veepee'],'French enterprise and deep tech'),
    # Asia continued (251-300)
    ('251','access-ventures','Access Ventures','Venture Capital','Seoul, Korea','$500M+','asia',['Korean enterprise tech'],'Korean enterprise technology'),
    ('252','b-dash','B Dash Ventures','Venture Capital','Tokyo, Japan','$200M+','asia',['Japanese enterprise startups'],'Japanese enterprise B2B'),
    ('253','cyber-agent','CyberAgent Ventures','Corporate VC','Tokyo, Japan','$1B+','asia',['Japanese digital media and enterprise'],'Japanese tech conglomerate CVC'),
    ('254','east-ventures','East Ventures','Venture Capital','Jakarta / Tokyo','$500M+','asia',['Tokopedia','Ruangguru','ShopBack'],'Southeast Asia tech leader'),
    ('255','gobi','Gobi Partners','Venture Capital','Kuala Lumpur / Shanghai','$1.5B+','asia',['ASEAN enterprise tech'],'Pan-Asian enterprise investing'),
    ('256','golden-gate','Golden Gate Ventures','Venture Capital','Singapore','$300M+','asia',['SE Asian fintech and enterprise'],'Singapore seed-stage pioneer'),
    ('257','infocomm','Infocomm Investments','Government Fund','Singapore','$500M+','asia',['Singapore digital economy'],'Singapore government tech fund'),
    ('258','jafco','JAFCO Group','Venture Capital','Tokyo, Japan','$3B+','asia',['Nomura NRI','SmartNews'],'Oldest Japanese VC — enterprise focus'),
    ('259','legends-capital','Legend Capital','Venture Capital','Beijing, China','$5B+','asia',['Lenovo legacy','Enterprise SaaS in China'],'Lenovo spinoff — Chinese enterprise'),
    ('260','matrix-india','Matrix Partners India','Venture Capital','Mumbai, India','$2B+','asia',['Razorpay','Ola','Country Delight'],'India B2B and fintech leader'),
    ('261','monks-hill','Monks Hill Ventures','Venture Capital','Singapore','$200M+','asia',['ASEAN enterprise B2B'],'SE Asia enterprise specialists'),
    ('262','openspace','OpenSpace Ventures','Venture Capital','Singapore','$500M+','asia',['Funding Societies','Carro'],'ASEAN Series A enterprise'),
    ('263','pavilion','Pavilion Capital','Sovereign Fund','Singapore','$2B+','asia',['Temasek subsidiary investments'],'Temasek subsidiary — Asia PE'),
    ('264','saison-capital','Saison Capital','Venture Capital','Singapore','$200M+','asia',['Credit Saison portfolio'],'Japan-Singapore enterprise bridge'),
    ('265','shimao','China Renaissance','Investment Bank/VC','Hong Kong/Beijing','$2B+','asia',['Didi','Meituan advisory'],'Chinese tech advisory and investing'),
    ('266','sirius','Sirius Venture Capital','Venture Capital','Melbourne, Australia','$200M+','asia',['Australian deep tech'],'Australian defense and enterprise tech'),
    ('267','strive','STRIVE','Venture Capital','Tokyo, Japan','$200M+','asia',['Japanese enterprise SaaS'],'Japan early-stage enterprise'),
    ('268','valar','Valar Ventures','Venture Capital','New York/London','$1B+','global',['TransferWise','N26','Xero'],'Peter Thiel — global fintech/enterprise'),
    ('269','venturra','Venturra','Venture Capital','Jakarta, Indonesia','$300M+','asia',['Indonesian enterprise tech'],'Indonesian Series A enterprise'),
    ('270','vy-capital','VY Capital','Growth Equity','Dubai / San Francisco','$5B+','middle-east',['SpaceX','Discord','Reddit'],'Crossover growth — UAE-US bridge'),
    # Additional global (271-300)
    ('271','kaszek','Kaszek Ventures','Venture Capital','Buenos Aires, Argentina','$3B+','global',['Nubank','Kavak','NotCo'],'LatAm technology leader'),
    ('272','monashees','Monashees','Venture Capital','São Paulo, Brazil','$2B+','global',['99','Neon','Loft'],'Brazilian enterprise tech'),
    ('273','valor-capital','Valor Capital Group','Venture Capital','NYC / São Paulo','$1B+','global',['Gympass','Loft','Creditas'],'US-Brazil bridge fund'),
    ('274','softbank-latam','SoftBank Latin America Fund','Venture Capital','Miami / São Paulo','$8B+','global',['Rappi','Kavak','Vtex'],'LatAm technology ecosystem'),
    ('275','flourish','Flourish Ventures','Impact VC','San Francisco / Bangalore','$1B+','global',['Chime','Chipper Cash','Andela'],'Fintech for prosperity — Omidyar'),
    ('276','catalyst-fund','Catalyst Fund','Impact VC','New York / Nairobi','$200M+','global',['African fintech startups'],'Emerging market fintech inclusion'),
    ('277','tlcom','TLcom Capital','Venture Capital','London / Lagos / Nairobi','$300M+','global',['African enterprise tech'],'Africa enterprise technology leader'),
    ('278','partech-africa','Partech Africa','Venture Capital','Dakar / Nairobi','$300M+','global',['African tech ecosystem'],'Pan-African enterprise investing'),
    ('279','sawari','Sawari Ventures','Venture Capital','Cairo, Egypt','$100M+','middle-east',['MENA tech startups'],'Egyptian enterprise technology'),
    ('280','novastar','Novastar Ventures','Venture Capital','Nairobi / London','$200M+','global',['East African enterprise'],'East Africa enterprise tech pioneer'),
    ('281','algebra','Algebra Ventures','Venture Capital','Cairo, Egypt','$100M+','middle-east',['Swvl','Halan','MaxAB'],'Egyptian enterprise and fintech'),
    ('282','cre-vc','CRE Venture Capital','Venture Capital','Hong Kong','$300M+','asia',['Asian fintech and enterprise'],'HK enterprise and compliance tech'),
    ('283','deep-knowledge','Deep Knowledge Ventures','Venture Capital','Hong Kong','$500M+','asia',['AI and longevity'],'AI governance and deep tech Asia'),
    ('284','dg-daiwa','DG Daiwa Ventures','Corporate VC','Tokyo, Japan','$300M+','asia',['Japanese enterprise fintech'],'Daiwa Securities CVC'),
    ('285','global-founders','Global Founders Capital','Venture Capital','Berlin / Singapore','$1B+','global',['Rocket Internet legacy'],'Transatlantic enterprise platforms'),
    ('286','horizons','Horizons Ventures','Venture Capital','Hong Kong','$2B+','asia',['DeepMind','Zoom','Spotify'],'Li Ka-shing — technology visionary'),
    ('287','james-riney','Coral Capital (James Riney)','Venture Capital','Tokyo, Japan','$300M+','asia',['SmartHR','inFlow'],'Japan enterprise SaaS leader'),
    ('288','kima','Kima Ventures','Venture Capital','Paris, France','$200M+','eu',['600+ startups globally'],'Xavier Niel — prolific angel/seed'),
    ('289','luxor','Luxor Capital','Hedge Fund / VC','New York, NY','$5B+','us',['Tech opportunistic investing'],'Crossover technology investing'),
    ('290','marathon','Marathon Asset Management','Asset Manager','New York, NY','$20B+','us',['Credit and technology'],'Credit-focused with tech allocation'),
    ('291','meritech','Meritech Capital','Growth Equity','Palo Alto, CA','$6B+','us',['Unity','Brex','Confluent'],'Enterprise growth specialist'),
    ('292','nerdy','Nerdy (Varsity Tutors parent)','Growth Equity','St. Louis, MO','$500M+','us',['EdTech enterprise'],'EdTech enterprise platform'),
    ('293','owl-rock','Blue Owl (ex-Owl Rock)','Private Credit','New York, NY','$170B+','us',['Enterprise tech lending'],'Technology-focused private credit'),
    ('294','permira','Permira','Private Equity','London / Menlo Park','$80B+','global',['Informatica','Genesys','Zendesk'],'Enterprise software buyout specialist'),
    ('295','providence','Providence Equity','Private Equity','Providence, RI','$50B+','us',['Blackboard','Houghton Mifflin'],'Education and media technology'),
    ('296','silver-lake','Silver Lake','Private Equity','Menlo Park / NYC','$100B+','global',['Dell','Airbnb','Unity'],'Enterprise tech PE leader'),
    ('297','TA-associates','TA Associates','Growth Equity','Boston, MA','$50B+','us',['RealPage','Calypso','Immuta'],'50+ year enterprise growth investor'),
    ('298','tw-telecom','TeleSoft Partners','Venture Capital','Saratoga, CA','$500M+','us',['Enterprise telecom and SaaS'],'Defense and enterprise telecom'),
    ('299','vitruvian','Vitruvian Partners','Growth Equity','London, UK','$10B+','eu',['Darktrace','Just Eat','Farfetch'],'European enterprise growth leader'),
    ('300','welsh-carson','Welsh Carson','Private Equity','New York, NY','$35B+','us',['Cotiviti','Optum360'],'Healthcare and enterprise IT PE'),
]

for f in _more_firms:
    INVESTOR_500_BATCH2.append({
        'num': int(f[0]), 'slug': f[1], 'name': f[2], 'type': f[3],
        'loc': f[4], 'aum': f[5], 'region': f[6], 
        'portfolio': f[7], 'thesis': f[8], 'web': ''
    })

# Enterprise targets (100 decks)
ENTERPRISE_TARGETS = []
_enterprise_list = [
    'AI Risk Management','Autonomous Vehicle Governance','Banking Compliance Platform',
    'Biotech Decision Intelligence','Board Governance Automation','Carbon Credit Verification',
    'Central Bank Digital Currency','Clinical Trial Decisions','Compliance Automation Suite',
    'Construction Risk Management','Contract Intelligence','Corporate Treasury AI',
    'Credit Risk AI Governance','Crisis Communication AI','Cross-Border Trade Compliance',
    'Cybersecurity Incident Response','Data Privacy Automation','Defense Acquisition AI',
    'Digital Twin Governance','Drug Discovery Decisions','E-Commerce Fraud Prevention',
    'Education Governance','Election Security','Emergency Response AI',
    'Energy Grid Optimization','Environmental Compliance','ESG Reporting Automation',
    'Export Control Compliance','FDA Submission AI','Financial Audit Automation',
    'Fleet Management AI','Food Safety Compliance','Fraud Detection Governance',
    'GDPR Compliance Engine','Healthcare AI Governance','HIPAA Audit Automation',
    'HR Decision Intelligence','Identity Verification','Immigration Case AI',
    'Infrastructure Monitoring','Insurance Claims AI','Intellectual Property AI',
    'Internal Audit AI','Investment Committee AI','IoT Security Governance',
    'KYC/AML Automation','Labor Law Compliance','Legal Discovery AI',
    'Loan Origination AI','Logistics Optimization','M&A Due Diligence AI',
    'Manufacturing Quality AI','Maritime Compliance','Medical Device Governance',
    'Military Intelligence AI','Mining Safety AI','Money Laundering Prevention',
    'Municipal Governance AI','Nuclear Safety AI','Oil & Gas Compliance',
    'Patent Analysis AI','Pension Fund Governance','Pharmaceutical Compliance',
    'Pipeline Safety AI','Police Oversight AI','Portfolio Risk AI',
    'Power Grid AI Governance','Precision Agriculture AI','Privacy Impact Assessment',
    'Procurement AI','Product Safety AI','Public Health AI',
    'Quality Assurance AI','Rail Safety AI','Real Estate Investment AI',
    'Regulatory Change Management','Renewable Energy AI','Research Ethics AI',
    'Retail Pricing AI','Risk Appetite Framework','Sanctions Screening AI',
    'School Board Governance','Securities Trading AI','Smart Building AI',
    'Social Media Governance','Space Operations AI','Sports Integrity AI',
    'Supply Chain AI','Tax Compliance AI','Telecom Spectrum AI',
    'Third-Party Risk AI','Trade Surveillance AI','Transportation Safety AI',
    'University Research AI','Utility Rate AI','Vendor Risk Management',
    'Veterinary AI Governance','Water Treatment AI','Wealth Management AI',
    'Workforce Planning AI','Zero Trust AI Security',
]
for i, name in enumerate(_enterprise_list[:100]):
    slug = name.lower().replace(' ', '-').replace('/', '-').replace('&', 'and')
    ENTERPRISE_TARGETS.append({'num': i+1, 'slug': slug, 'name': name})

# Regional targets (122 decks)
REGIONAL_TARGETS = []
_regional_list = [
    'United States - Northeast','United States - Southeast','United States - Midwest',
    'United States - West Coast','United States - Texas','United States - DC/Virginia',
    'Canada - Ontario','Canada - Quebec','Canada - British Columbia',
    'United Kingdom - London','United Kingdom - Scotland','United Kingdom - Manchester',
    'Germany - Berlin','Germany - Munich','Germany - Frankfurt',
    'France - Paris','France - Lyon','France - Toulouse',
    'Switzerland - Zurich','Switzerland - Geneva','Netherlands - Amsterdam',
    'Belgium - Brussels','Luxembourg','Sweden - Stockholm',
    'Norway - Oslo','Denmark - Copenhagen','Finland - Helsinki',
    'Ireland - Dublin','Spain - Madrid','Spain - Barcelona',
    'Portugal - Lisbon','Italy - Milan','Italy - Rome',
    'Austria - Vienna','Poland - Warsaw','Czech Republic - Prague',
    'Estonia - Tallinn','Latvia - Riga','Lithuania - Vilnius',
    'Romania - Bucharest','Greece - Athens','Israel - Tel Aviv',
    'Israel - Jerusalem','Turkey - Istanbul','UAE - Dubai',
    'UAE - Abu Dhabi','Saudi Arabia - Riyadh','Saudi Arabia - NEOM',
    'Qatar - Doha','Bahrain - Manama','Kuwait - Kuwait City',
    'Oman - Muscat','Jordan - Amman','Egypt - Cairo',
    'South Africa - Johannesburg','South Africa - Cape Town','Nigeria - Lagos',
    'Kenya - Nairobi','Ghana - Accra','Rwanda - Kigali',
    'Morocco - Casablanca','Ethiopia - Addis Ababa','Tanzania - Dar es Salaam',
    'Japan - Tokyo','Japan - Osaka','South Korea - Seoul',
    'China - Beijing','China - Shanghai','China - Shenzhen',
    'Hong Kong','Taiwan - Taipei','Singapore',
    'India - Mumbai','India - Bangalore','India - Delhi',
    'India - Hyderabad','Indonesia - Jakarta','Malaysia - Kuala Lumpur',
    'Thailand - Bangkok','Vietnam - Ho Chi Minh City','Philippines - Manila',
    'Australia - Sydney','Australia - Melbourne','New Zealand - Auckland',
    'Brazil - São Paulo','Brazil - Rio de Janeiro','Mexico - Mexico City',
    'Mexico - Monterrey','Colombia - Bogotá','Argentina - Buenos Aires',
    'Chile - Santiago','Peru - Lima','Costa Rica - San José',
    'Panama - Panama City','Uruguay - Montevideo','Ecuador - Quito',
    'EU - Brussels (EU AI Act)','EU - Strasbourg (Regulation)',
    'EU - Frankfurt (ECB/Finance)','EU - Luxembourg (EIB/ESM)',
    'NATO - Brussels (Defense)','UN - New York (SDGs)',
    'UN - Geneva (UNCTAD)','World Bank - Washington DC',
    'IMF - Washington DC','OECD - Paris',
    'WEF - Davos','ASEAN - Jakarta',
    'African Union - Addis Ababa','GCC - Riyadh',
    'APEC - Singapore','G7 Rotating',
    'G20 Rotating','BRICS - Johannesburg',
    'Pacific Islands Forum - Suva','Caribbean Community - Georgetown',
    'Nordic Council - Copenhagen','Mercosur - Montevideo',
    'SCO - Beijing','Arab League - Cairo',
]
for i, name in enumerate(_regional_list[:122]):
    slug = name.lower().replace(' ', '-').replace('/', '-').replace('(', '').replace(')', '').replace(',', '')
    slug = slug.replace('---', '-').replace('--', '-')
    REGIONAL_TARGETS.append({'num': i+1, 'slug': slug, 'name': name})

# Investor-50 extra themes (49 more)
INVESTOR_50_EXTRA = []
_themes = [
    'AI Agent Orchestration','Autonomous Systems Governance','Blockchain Governance',
    'Board Decision Automation','Carbon Market Intelligence','Central Bank AI',
    'Clinical AI Governance','Compliance Cost Reduction','Conversational AI Governance',
    'Crypto Compliance','Cyber Insurance','Data Mesh Governance',
    'Decision Debt Reduction','Defense AI Ethics','Digital Identity Governance',
    'Edge AI Governance','Embedded Finance AI','Energy Transition AI',
    'EU Digital Markets Act','Evidence-Based Policy','Explainable AI',
    'Federal AI Mandate','Financial Crime AI','Generative AI Governance',
    'Global AI Standards','Green Finance AI','Healthcare AI Ethics',
    'Human-AI Collaboration','Institutional Knowledge AI','Insurance AI Risk',
    'IoT Decision Intelligence','Judicial AI Governance','Knowledge Graph Governance',
    'Large Language Model Governance','Legal AI Compliance','Machine Learning Ops Governance',
    'Maritime AI Governance','Model Risk Management','National AI Strategy',
    'Nuclear AI Safety','Open Banking AI','Pension AI Governance',
    'Pharmaceutical AI Compliance','Post-Merger AI Integration','Quantum AI Governance',
    'Real-Time Compliance','Responsible AI Framework','Supply Chain AI Ethics',
    'Trade Finance AI',
]
for i, name in enumerate(_themes[:49]):
    slug = name.lower().replace(' ', '-').replace('/', '-')
    INVESTOR_50_EXTRA.append({'num': 52+i, 'slug': slug, 'name': name})

# Generic investor extra (5 more)
INVESTOR_EXTRA = [
    {'num': 11, 'slug': 'eu-ai-act-deep-dive', 'name': 'EU AI Act Deep Dive'},
    {'num': 12, 'slug': 'dora-compliance', 'name': 'DORA Financial Compliance'},
    {'num': 13, 'slug': 'defense-cleared', 'name': 'Defense / Cleared Markets'},
    {'num': 14, 'slug': 'open-source-moat', 'name': 'Open Source Strategy & Moat'},
    {'num': 15, 'slug': 'ai-liability-insurance', 'name': 'AI Liability & Insurance'},
]

# Gamma extra (4 more)
GAMMA_EXTRA = [
    {'num': 7, 'slug': 'healthcare-vertical', 'name': 'Healthcare Vertical'},
    {'num': 8, 'slug': 'financial-compliance', 'name': 'Financial Compliance'},
    {'num': 9, 'slug': 'ai-governance-market', 'name': 'AI Governance Market'},
    {'num': 10, 'slug': 'open-source-strategy', 'name': 'Open Source Strategy'},
]


def main():
    test_mode = '--test' in sys.argv
    category_filter = None
    for arg in sys.argv[1:]:
        if arg.startswith('--category='):
            category_filter = arg.split('=')[1]
    
    generated = 0
    
    # ── Investor-500 (firms 121-300) ──
    if not category_filter or category_filter == 'investor-500':
        all_investors = INVESTOR_500_NEW + INVESTOR_500_BATCH2
        output_dir = os.path.join(BASE, 'investor-500', 'pptx')
        print(f'\nGenerating {len(all_investors)} investor-500 decks...')
        for firm in (all_investors[:1] if test_mode else all_investors):
            path = generate_investor_deck(firm, output_dir)
            embed_screenshots(path)
            generated += 1
            if generated % 20 == 0:
                print(f'  {generated} generated...')
        print(f'  investor-500: {min(len(all_investors), 1 if test_mode else len(all_investors))} decks')
    
    # ── Sales extra (26-75) ──
    if not category_filter or category_filter == 'sales':
        output_dir = os.path.join(BASE, 'sales', 'pptx')
        print(f'\nGenerating {len(SALES_EXTRA)} extra sales decks...')
        for target in (SALES_EXTRA[:1] if test_mode else SALES_EXTRA):
            path = generate_sales_deck(target, output_dir)
            embed_screenshots(path)
            generated += 1
        print(f'  sales extra: {min(len(SALES_EXTRA), 1 if test_mode else len(SALES_EXTRA))} decks')
    
    # ── Design Partner extra (26-75) ──
    if not category_filter or category_filter == 'design-partner':
        output_dir = os.path.join(BASE, 'design-partner', 'pptx')
        print(f'\nGenerating {len(DESIGN_PARTNER_EXTRA)} extra design-partner decks...')
        for target in (DESIGN_PARTNER_EXTRA[:1] if test_mode else DESIGN_PARTNER_EXTRA):
            path = generate_design_partner_deck(target, output_dir)
            embed_screenshots(path)
            generated += 1
        print(f'  design-partner extra: {min(len(DESIGN_PARTNER_EXTRA), 1 if test_mode else len(DESIGN_PARTNER_EXTRA))} decks')
    
    # ── Enterprise (1-100) ──
    if not category_filter or category_filter == 'enterprise':
        output_dir = os.path.join(BASE, 'enterprise', 'pptx')
        print(f'\nGenerating {len(ENTERPRISE_TARGETS)} enterprise decks...')
        for target in (ENTERPRISE_TARGETS[:1] if test_mode else ENTERPRISE_TARGETS):
            path = generate_thematic_deck(target, output_dir)
            embed_screenshots(path)
            generated += 1
        print(f'  enterprise: {min(len(ENTERPRISE_TARGETS), 1 if test_mode else len(ENTERPRISE_TARGETS))} decks')
    
    # ── Regional (1-122) ──
    if not category_filter or category_filter == 'regional':
        output_dir = os.path.join(BASE, 'regional', 'pptx')
        print(f'\nGenerating {len(REGIONAL_TARGETS)} regional decks...')
        for target in (REGIONAL_TARGETS[:1] if test_mode else REGIONAL_TARGETS):
            path = generate_thematic_deck(target, output_dir)
            embed_screenshots(path)
            generated += 1
        print(f'  regional: {min(len(REGIONAL_TARGETS), 1 if test_mode else len(REGIONAL_TARGETS))} decks')
    
    # ── Investor-50 extra (52-100) ──
    if not category_filter or category_filter == 'investor-50':
        output_dir = os.path.join(BASE, 'investor-50', 'pptx')
        print(f'\nGenerating {len(INVESTOR_50_EXTRA)} extra investor-50 decks...')
        for target in (INVESTOR_50_EXTRA[:1] if test_mode else INVESTOR_50_EXTRA):
            path = generate_thematic_deck(target, output_dir)
            embed_screenshots(path)
            generated += 1
        print(f'  investor-50 extra: {min(len(INVESTOR_50_EXTRA), 1 if test_mode else len(INVESTOR_50_EXTRA))} decks')
    
    # ── Investor extra (11-15) ──
    if not category_filter or category_filter == 'investor':
        output_dir = os.path.join(BASE, 'investor', 'pptx')
        print(f'\nGenerating {len(INVESTOR_EXTRA)} extra investor decks...')
        for target in (INVESTOR_EXTRA[:1] if test_mode else INVESTOR_EXTRA):
            path = generate_thematic_deck(target, output_dir, 'investor-50')
            embed_screenshots(path)
            generated += 1
        print(f'  investor extra: {min(len(INVESTOR_EXTRA), 1 if test_mode else len(INVESTOR_EXTRA))} decks')
    
    # ── Gamma extra (7-10) ──
    if not category_filter or category_filter == 'gamma':
        output_dir = os.path.join(BASE, 'gamma', 'pptx')
        print(f'\nGenerating {len(GAMMA_EXTRA)} extra gamma decks...')
        for target in (GAMMA_EXTRA[:1] if test_mode else GAMMA_EXTRA):
            path = generate_thematic_deck(target, output_dir, 'gamma')
            embed_screenshots(path)
            generated += 1
        print(f'  gamma extra: {min(len(GAMMA_EXTRA), 1 if test_mode else len(GAMMA_EXTRA))} decks')
    
    print(f'\nDONE: {generated} new decks generated')
    
    # Count total
    total = 0
    for folder in ['pptx', 'investor/pptx', 'investor-50/pptx', 'investor-100/pptx', 
                    'investor-500/pptx', 'sales/pptx', 'design-partner/pptx', 
                    'gamma/pptx', 'enterprise/pptx', 'regional/pptx']:
        path = os.path.join(BASE, folder)
        if os.path.exists(path):
            count = len(glob.glob(os.path.join(path, '*.pptx')))
            total += count
            if count > 0:
                print(f'  {folder}: {count}')
    print(f'  GRAND TOTAL PPTX: {total}')


if __name__ == '__main__':
    main()
