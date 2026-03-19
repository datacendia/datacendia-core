"""
Audit all PPTX files for Picture9.jpg references and broken images.
This will identify where the Picture9.jpg issue exists.
"""
import os, sys, glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = os.path.dirname(os.path.abspath(__file__))
PPTX_DIRS = [
    os.path.join(BASE, 'investor-100', 'pptx'),
    os.path.join(BASE, 'investor-500', 'pptx'),
    os.path.join(BASE, 'sales', 'pptx'),
    os.path.join(BASE, 'design-partner', 'pptx'),
]

def audit_deck(pptx_path):
    """Audit a single PPTX for Picture9 references and broken images."""
    issues = []
    name = os.path.basename(pptx_path)
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return [(name, "ERROR", f"Cannot open: {e}")]
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        
        # Check all shapes
        for shape in slide.shapes:
            # Check for picture shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    size = len(img.blob)
                    if size < 1000:  # Very small = likely broken/placeholder
                        issues.append((name, slide_num, f"Small image ({size} bytes)"))
                except Exception as e:
                    issues.append((name, slide_num, f"BROKEN image: {str(e)[:50]}"))
            
            # Check text references to Picture9
            if shape.has_text_frame:
                text = shape.text_frame.text
                if 'Picture9' in text or 'picture9' in text.lower():
                    issues.append((name, slide_num, f"Text ref to Picture9: {text[:80]}"))
            
            # Check for image placeholders or broken refs
            if hasattr(shape, 'image') and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                issues.append((name, slide_num, f"Non-picture with image attr (type={shape.shape_type})"))
    
    return issues

def main():
    all_issues = []
    
    for pptx_dir in PPTX_DIRS:
        if not os.path.isdir(pptx_dir):
            continue
        
        pptx_files = glob.glob(os.path.join(pptx_dir, '*.pptx'))
        print(f"\nAuditing {len(pptx_files)} files in {pptx_dir}")
        
        for pptx_file in pptx_files:
            issues = audit_deck(pptx_file)
            all_issues.extend(issues)
    
    print(f"\n{'='*60}")
    print(f"PICTURE9 AUDIT RESULTS")
    print(f"{'='*60}")
    
    if not all_issues:
        print("No Picture9.jpg references or broken images found in any PPTX files.")
        return
    
    print(f"Found {len(all_issues)} issues:")
    for filename, slide, issue in all_issues:
        print(f"  {filename} - Slide {slide}: {issue}")
    
    # Group by type
    broken_images = [i for i in all_issues if 'BROKEN' in i[2] or 'Small image' in i[2]]
    text_refs = [i for i in all_issues if 'Text ref to Picture9' in i[2]]
    
    print(f"\nSummary:")
    print(f"  Broken/Small images: {len(broken_images)}")
    print(f"  Text references to Picture9: {len(text_refs)}")

if __name__ == '__main__':
    main()
