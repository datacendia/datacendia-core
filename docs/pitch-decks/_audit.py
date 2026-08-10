"""
Full audit of every PPTX file: extract all text from every slide,
identify existing images, determine target audience.
Output: structured JSON for planning proper screenshot placement.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
import os, glob, json

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

def audit_deck(pptx_path):
    """Extract complete structure of a PPTX file."""
    prs = Presentation(pptx_path)
    deck = {
        'file': os.path.relpath(pptx_path, BASE),
        'slide_width_inches': round(prs.slide_width / 914400, 2),
        'slide_height_inches': round(prs.slide_height / 914400, 2),
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    for i, slide in enumerate(prs.slides):
        slide_data = {
            'number': i + 1,
            'texts': [],
            'images': [],
            'all_text_combined': ''
        }
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data['images'].append({
                    'width_in': round(shape.width / 914400, 2),
                    'height_in': round(shape.height / 914400, 2),
                    'left_in': round(shape.left / 914400, 2),
                    'top_in': round(shape.top / 914400, 2),
                })
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_data['texts'].append(t)
        
        slide_data['all_text_combined'] = ' | '.join(slide_data['texts'])
        deck['slides'].append(slide_data)
    
    return deck

# Audit ALL decks
folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 
           'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

all_decks = []
for folder in folders:
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    for f in files:
        try:
            deck = audit_deck(f)
            all_decks.append(deck)
        except Exception as e:
            print(f'ERROR: {f}: {e}')

# Save full audit
output = os.path.join(BASE, '_audit_results.json')
with open(output, 'w', encoding='utf-8') as fp:
    json.dump(all_decks, fp, indent=2, ensure_ascii=False)

# Print summary
print(f'Audited {len(all_decks)} decks')
print(f'Output: {output}')

# Print category summary
for folder in folders:
    decks_in = [d for d in all_decks if d['file'].startswith(folder.replace('/', '\\'))]
    if not decks_in:
        continue
    print(f'\n=== {folder} ({len(decks_in)} decks) ===')
    for d in decks_in[:2]:
        print(f'  {d["file"]} — {d["slide_count"]} slides, {d["slide_width_inches"]}x{d["slide_height_inches"]}in')
        for s in d['slides']:
            img_count = len(s['images'])
            img_str = f' [{img_count} img]' if img_count else ''
            # Get slide title (first meaningful text that isn't DATACENDIA or copyright)
            title = '(empty)'
            for t in s['texts']:
                if t not in ('DATACENDIA',) and 'datacendia.com' not in t.lower() and '2024' not in t:
                    title = t[:80]
                    break
            print(f'    S{s["number"]}{img_str}: {title}')
