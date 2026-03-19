"""
Find duplicate pitch decks: same target company appearing in multiple categories.
Also find PDFs that ended up in wrong folders (pptx/ instead of pdf/).
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os, glob, hashlib
from collections import defaultdict

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

# 1. Find PDFs accidentally saved in pptx/ folders
print("=" * 60)
print("PDFs in wrong folders (pptx/ dirs)")
print("=" * 60)
wrong_pdfs = []
for root, dirs, files in os.walk(BASE):
    for f in files:
        full = os.path.join(root, f)
        rel = os.path.relpath(full, BASE)
        if f.endswith('.pdf') and '\\pptx\\' in full:
            wrong_pdfs.append(full)
            print(f"  {rel}")
print(f"\nTotal misplaced PDFs: {len(wrong_pdfs)}")

# 2. Extract target names from ALL decks to find cross-category duplicates
print("\n" + "=" * 60)
print("Cross-category duplicate targets")
print("=" * 60)

target_map = defaultdict(list)  # normalized_target -> [(category, filename)]

folders = {
    'pptx': 'main',
    'investor/pptx': 'investor',
    'investor-50/pptx': 'investor-50',
    'investor-100/pptx': 'investor-100',
    'sales/pptx': 'sales',
    'design-partner/pptx': 'design-partner',
    'gamma/pptx': 'gamma',
}

for folder, category in folders.items():
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    for f in sorted(glob.glob(os.path.join(path, '*.pptx'))):
        name = os.path.basename(f)
        try:
            prs = Presentation(f)
            # Extract target from first slide text
            target = None
            for slide in prs.slides[:2]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if 'prepared for' in t.lower():
                                target = t.replace('PREPARED FOR ', '').strip()
                                break
                    if target:
                        break
                if target:
                    break
            
            if not target:
                # Use filename as target
                target = name.replace('.pptx', '').split('-', 1)[-1] if '-' in name else name
            
            normalized = target.lower().strip()
            target_map[normalized].append((category, name, target))
        except:
            pass

# Find targets that appear in multiple categories
dupes_found = 0
for norm_target, entries in sorted(target_map.items()):
    categories = set(e[0] for e in entries)
    if len(categories) > 1:
        dupes_found += 1
        print(f"\n  TARGET: {entries[0][2]}")
        for cat, fname, tgt in entries:
            print(f"    [{cat}] {fname}")

if dupes_found == 0:
    print("  None found - all targets are unique per category")

# 3. Check for exact file duplicates (same hash)
print("\n" + "=" * 60)
print("Exact file duplicates (identical content)")
print("=" * 60)

hash_map = defaultdict(list)
for root, dirs, files in os.walk(BASE):
    for f in files:
        if f.endswith(('.pptx', '.pdf')):
            full = os.path.join(root, f)
            h = hashlib.md5(open(full, 'rb').read()).hexdigest()
            rel = os.path.relpath(full, BASE)
            hash_map[h].append(rel)

exact_dupes = 0
for h, paths in hash_map.items():
    if len(paths) > 1:
        # Filter out legitimate pptx/pdf pairs
        extensions = set(os.path.splitext(p)[1] for p in paths)
        if extensions == {'.pptx', '.pdf'}:
            continue  # Expected - pptx and its pdf
        exact_dupes += 1
        print(f"\n  Hash: {h[:12]}...")
        for p in paths:
            print(f"    {p}")

if exact_dupes == 0:
    print("  None found - no identical files")

# 4. Summary
print("\n" + "=" * 60)
print("SUMMARY")
print("=" * 60)
total_pptx = len(glob.glob(os.path.join(BASE, '**/*.pptx'), recursive=True))
total_pdf = len(glob.glob(os.path.join(BASE, '**/*.pdf'), recursive=True))
print(f"Total PPTX: {total_pptx}")
print(f"Total PDF: {total_pdf}")
print(f"Misplaced PDFs (in pptx/ folders): {len(wrong_pdfs)}")
print(f"Cross-category duplicates: {dupes_found}")
print(f"Exact file duplicates: {exact_dupes}")
