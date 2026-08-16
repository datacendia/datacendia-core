import os, glob
from pptx import Presentation

BASE = r'c:\Users\Stu\datacendia-core\docs\pitch-decks'
checks = [
    ('design-partner/pptx', 'Clinical decision support with HIPAA-compliant audit trails'),
    ('design-partner/pptx', 'patient safety as the prime directive'),
    ('enterprise/pptx', 'AI Governance for Regulated Financial Institutions'),
    ('enterprise/pptx', '& Capital Markets'),
    ('enterprise/pptx', 'DORA'),
    ('regional/pptx', 'AI Governance for Regulated Financial Institutions'),
    ('regional/pptx', '& Capital Markets'),
    ('regional/pptx', 'DORA'),
    ('gamma/pptx', 'Decision CrisisImmunizationInfrastructure'),
]
for folder, needle in checks:
    path = os.path.join(BASE, folder)
    hits = 0
    total = 0
    examples = []
    for f in sorted(glob.glob(os.path.join(path, '*.pptx'))):
        total += 1
        try:
            prs = Presentation(f)
            text = '\n'.join(
                para.text.strip()
                for slide in prs.slides
                for shape in slide.shapes if shape.has_text_frame
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
            if needle in text:
                hits += 1
                if len(examples) < 5:
                    examples.append(os.path.basename(f))
        except Exception:
            pass
    print(f'{folder} | {needle} | {hits}/{total} | {", ".join(examples)}')
