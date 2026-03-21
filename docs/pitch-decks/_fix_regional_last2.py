"""Fix the last 2 regional decks (121, 122) that still have banking template bleed."""
import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

REPLACEMENTS = [
    ('& Capital Markets', ''),
    ('AI Governance for Regulated Financial Institutions', 'Decision Evidence Infrastructure for AI Systems'),
    ('DORA', 'EU AI Act'),
    ('Jan 2025 \u2014 AI Operational Resilience', '2025\u20132026 \u2014 High-Risk AI Governance Mandated'),
    ('Basel III', 'ISO/IEC 42001'),
    ('Auditable Risk Models Required', 'AI Management System Certification'),
    ('MiFID II', 'NIST AI RMF'),
    ('Full Audit Trail for Algo Trading', 'Risk Management for Trustworthy AI'),
    ('Banks Deploy AI Without Audit Trails', 'Organizations Deploy AI Without Decision Evidence'),
    ('Credit decisions influenced by AI with no explainability for regulators',
     'High-stakes decisions influenced by AI with no explainability for stakeholders'),
    ('DORA requires operational resilience documentation for all AI systems by Jan 2025',
     'EU AI Act requires risk management and human oversight for high-risk AI systems'),
    ('Basel III mandates auditable risk models \u2014 AI black boxes fail this test',
     'ISO/IEC 42001 mandates AI management systems \u2014 black-box AI fails this test'),
    ('MiFID II requires full audit trails for algorithmic trading decisions',
     'Regulators and boards increasingly require full audit trails for AI-influenced decisions'),
    ('credit, trading, and risk decision', 'high-stakes decision'),
    ('Credit decision review with dissent tracking', 'Decision review with dissent tracking'),
    ('M&A due diligence with full evidence trail', 'Due diligence with full evidence trail'),
    ('DORA compliance with automated evidence', 'Compliance with automated evidence generation'),
    ('Investment committee structured deliberation', 'Committee-level structured deliberation'),
    ('Regulatory change mapping to P&L impact', 'Regulatory change mapping to operational impact'),
    ('Banking is the largest regulated vertical. DORA enforcement creates immediate demand. Only platform providing cryptographic decision evidence \u2014 a regulatory requirement with zero incumbent solutions.',
     'AI governance is becoming mandatory across regulated industries. EU AI Act enforcement creates immediate demand. Only platform providing cryptographic decision evidence \u2014 a regulatory requirement with zero incumbent solutions.'),
    ("Built the exact data infrastructure that financial institutions need for AI governance \u2014 at one of the world's largest asset managers.",
     'Built the data infrastructure and decision workflows required for regulated, high-stakes AI deployments at institutional scale.'),
]

fixed = 0
for fn in ['121-east-african-community.pptx', '122-sadc-gaborone.pptx']:
    fp = os.path.join(BASE, 'regional', 'pptx', fn)
    if not os.path.exists(fp):
        print(f'  SKIP (not found): {fn}')
        continue
    prs = Presentation(fp)
    ch = False
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for pa in sh.text_frame.paragraphs:
                for r in pa.runs:
                    for old, new in REPLACEMENTS:
                        if old in r.text:
                            r.text = r.text.replace(old, new)
                            ch = True
    if ch:
        prs.save(fp)
        fixed += 1
        print(f'  Fixed {fn}')
    else:
        print(f'  No changes needed: {fn}')

print(f'Regional fixed: {fixed}/2')
