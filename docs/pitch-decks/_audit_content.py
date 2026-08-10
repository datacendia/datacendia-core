"""Extract all text from the a16z deck and other templates to audit content."""
import os
from pptx import Presentation

BASE = os.path.dirname(__file__)
templates = [
    ('investor-100/pptx/001-a16z.pptx', 'INVESTOR (a16z)'),
    ('sales/pptx/01-fortune-500.pptx', 'SALES (Fortune 500)'),
    ('design-partner/pptx/01-healthcare.pptx', 'DESIGN PARTNER (Healthcare)'),
]

for rel, label in templates:
    path = os.path.join(BASE, rel)
    if not os.path.exists(path):
        print(f'\n*** {label}: FILE NOT FOUND ***')
        continue
    prs = Presentation(path)
    print(f'\n{"="*80}')
    print(f'{label} — {rel}')
    print(f'{"="*80}')
    for si, slide in enumerate(prs.slides):
        print(f'\n--- Slide {si+1} ---')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        print(f'  [{shape.shape_type}] {text}')
            if shape.shape_type == 13:  # Picture
                print(f'  [IMAGE] {shape.width}x{shape.height} at ({shape.left},{shape.top})')
