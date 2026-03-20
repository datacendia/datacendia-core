"""Verify screenshots are still embedded in all pitch decks after text fixes."""
import os, sys, io, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = os.path.dirname(os.path.abspath(__file__))
folders = ['pptx','investor/pptx','investor-50/pptx','investor-100/pptx',
           'investor-500/pptx','sales/pptx','design-partner/pptx',
           'gamma/pptx','enterprise/pptx','regional/pptx']

total = 0
with_images = 0
without_images = 0
missing_list = []

for folder in folders:
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    for f in files:
        total += 1
        try:
            prs = Presentation(f)
            has_img = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        has_img = True
                        break
                if has_img:
                    break
            if has_img:
                with_images += 1
            else:
                without_images += 1
                missing_list.append(os.path.relpath(f, BASE))
        except Exception as e:
            print(f'ERROR: {f}: {e}')

print(f'Total PPTX scanned: {total}')
print(f'With screenshots:   {with_images}')
print(f'WITHOUT screenshots: {without_images}')
if missing_list:
    print(f'\nDecks missing screenshots:')
    for m in missing_list[:20]:
        print(f'  {m}')
    if len(missing_list) > 20:
        print(f'  ... and {len(missing_list)-20} more')
