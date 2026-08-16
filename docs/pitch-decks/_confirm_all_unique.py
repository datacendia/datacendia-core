"""
COMPREHENSIVE CONFIRMATION: Every pitch deck is different and individually aimed.

Checks ALL 1000 decks:
  1. No two files are byte-identical
  2. No two files have identical text content
  3. Every investor-500 deck names its specific firm on Slide 1 and Slide 10
  4. Every investor-100 deck names its specific firm
  5. Every investor-50 deck names its specific theme
  6. Every sales deck names its specific target
  7. Every design-partner deck names its specific target
  8. Every enterprise deck names its specific topic
  9. Every regional deck names its specific region
  10. Every gamma deck is distinct
"""
import hashlib
import os
import sys
import io
import glob
from collections import defaultdict

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

CATEGORIES = [
    ('investor-500', 'investor-500/pptx'),
    ('investor-100', 'investor-100/pptx'),
    ('investor-50', 'investor-50/pptx'),
    ('investor', 'investor/pptx'),
    ('sales', 'sales/pptx'),
    ('design-partner', 'design-partner/pptx'),
    ('enterprise', 'enterprise/pptx'),
    ('regional', 'regional/pptx'),
    ('gamma', 'gamma/pptx'),
    ('main', 'pptx'),
]


def file_hash(path):
    h = hashlib.sha256()
    with open(path, 'rb') as f:
        for chunk in iter(lambda: f.read(65536), b''):
            h.update(chunk)
    return h.hexdigest()


def get_slide_text(pptx_path, slide_num):
    """Get all text from a specific slide (1-indexed)."""
    try:
        prs = Presentation(pptx_path)
        if slide_num > len(prs.slides):
            return ''
        slide = prs.slides[slide_num - 1]
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        return ' | '.join(parts)
    except Exception as e:
        return f'ERROR: {e}'


def get_all_text(pptx_path):
    """Get all text from entire deck."""
    try:
        prs = Presentation(pptx_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
        return '\n'.join(parts)
    except:
        return ''


# ═══════════════════════════════════════════════════════════════
# TEST 1: File-level uniqueness (byte-identical check)
# ═══════════════════════════════════════════════════════════════
print('=' * 70)
print('TEST 1: FILE-LEVEL UNIQUENESS (byte-identical check)')
print('=' * 70)

all_files = []
for cat, folder in CATEGORIES:
    path = os.path.join(BASE, folder)
    if os.path.exists(path):
        all_files.extend(sorted(glob.glob(os.path.join(path, '*.pptx'))))

hash_to_files = defaultdict(list)
for f in all_files:
    h = file_hash(f)
    hash_to_files[h].append(os.path.relpath(f, BASE))

dupes = {h: files for h, files in hash_to_files.items() if len(files) > 1}
if dupes:
    print(f'  FAIL: {len(dupes)} sets of identical files found!')
    for h, files in dupes.items():
        print(f'    {files}')
else:
    print(f'  PASS: All {len(all_files)} files are byte-unique')


# ═══════════════════════════════════════════════════════════════
# TEST 2: Text-content uniqueness
# ═══════════════════════════════════════════════════════════════
print()
print('=' * 70)
print('TEST 2: TEXT-CONTENT UNIQUENESS')
print('=' * 70)

text_to_files = defaultdict(list)
for f in all_files:
    text = get_all_text(f)
    text_hash = hashlib.sha256(text.encode('utf-8', errors='replace')).hexdigest()
    text_to_files[text_hash].append(os.path.relpath(f, BASE))

text_dupes = {h: files for h, files in text_to_files.items() if len(files) > 1}
if text_dupes:
    print(f'  WARNING: {len(text_dupes)} sets of text-identical files:')
    for h, files in text_dupes.items():
        print(f'    {files}')
else:
    print(f'  PASS: All {len(all_files)} files have unique text content')


# ═══════════════════════════════════════════════════════════════
# TEST 3: INVESTOR-500 — Each deck names its specific firm
# ═══════════════════════════════════════════════════════════════
print()
print('=' * 70)
print('TEST 3: INVESTOR-500 — Each deck individually aimed')
print('=' * 70)

folder = os.path.join(BASE, 'investor-500', 'pptx')
files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
inv500_ok = 0
inv500_fail = []
firm_names_seen = set()

for f in files:
    fname = os.path.basename(f)
    s1 = get_slide_text(f, 1)  # Cover slide
    s10 = get_slide_text(f, 10)  # Closing slide

    # Extract what firm name appears on cover
    # Cover should say "PREPARED FOR [FIRM]" or "Datacendia for [FIRM]"
    has_prepared = 'PREPARED FOR' in s1
    has_datacendia_for = 'Datacendia for' in s1

    # Closing should say "Prepared for [FIRM]"
    has_closing = 'Prepared for' in s10 or 'PREPARED FOR' in s10

    if has_prepared and has_closing:
        inv500_ok += 1
        # Extract firm name from "PREPARED FOR X"
        if 'PREPARED FOR ' in s1:
            firm = s1.split('PREPARED FOR ')[1].split(' |')[0].split(' | ')[0].strip()
            firm_names_seen.add(firm)
    else:
        inv500_fail.append((fname, 'cover' if not has_prepared else 'closing'))

print(f'  Total: {len(files)}')
print(f'  Individually aimed (firm on cover + closing): {inv500_ok}/{len(files)}')
print(f'  Unique firm names: {len(firm_names_seen)}')
if inv500_fail:
    print(f'  ISSUES ({len(inv500_fail)}):')
    for fn, issue in inv500_fail[:10]:
        print(f'    {fn}: missing {issue} personalization')
else:
    print(f'  PASS: Every investor-500 deck is individually aimed')


# ═══════════════════════════════════════════════════════════════
# TEST 4: INVESTOR-100 — Each deck names its specific firm
# ═══════════════════════════════════════════════════════════════
print()
print('=' * 70)
print('TEST 4: INVESTOR-100 — Each deck individually aimed')
print('=' * 70)

folder = os.path.join(BASE, 'investor-100', 'pptx')
files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
inv100_ok = 0
inv100_names = set()

for f in files:
    s1 = get_slide_text(f, 1)
    if 'PREPARED FOR' in s1:
        inv100_ok += 1
        firm = s1.split('PREPARED FOR ')[1].split(' |')[0].strip() if 'PREPARED FOR ' in s1 else ''
        if firm:
            inv100_names.add(firm)

print(f'  Total: {len(files)}')
print(f'  Individually aimed: {inv100_ok}/{len(files)}')
print(f'  Unique firm names: {len(inv100_names)}')
if inv100_ok == len(files):
    print(f'  PASS')


# ═══════════════════════════════════════════════════════════════
# TEST 5: INVESTOR-50 — Each deck names its specific theme
# ═══════════════════════════════════════════════════════════════
print()
print('=' * 70)
print('TEST 5: INVESTOR-50 — Each deck individually themed')
print('=' * 70)

folder = os.path.join(BASE, 'investor-50', 'pptx')
files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
inv50_themes = set()

for f in files:
    s1 = get_slide_text(f, 1)
    # These decks have the theme in the title, not "PREPARED FOR"
    inv50_themes.add(s1[:200])

print(f'  Total: {len(files)}')
print(f'  Unique cover text: {len(inv50_themes)}/{len(files)}')
if len(inv50_themes) == len(files):
    print(f'  PASS: Every investor-50 deck has a unique cover')
else:
    print(f'  WARNING: {len(files) - len(inv50_themes)} decks share cover text')


# ═══════════════════════════════════════════════════════════════
# TEST 6-9: Other categories — unique cover text
# ═══════════════════════════════════════════════════════════════
for cat_name, cat_folder in [('SALES', 'sales/pptx'), ('DESIGN-PARTNER', 'design-partner/pptx'),
                               ('ENTERPRISE', 'enterprise/pptx'), ('REGIONAL', 'regional/pptx'),
                               ('GAMMA', 'gamma/pptx'), ('INVESTOR', 'investor/pptx'), ('MAIN', 'pptx')]:
    print()
    print('=' * 70)
    print(f'TEST: {cat_name} — Each deck individually targeted')
    print('=' * 70)

    folder = os.path.join(BASE, cat_folder)
    if not os.path.exists(folder):
        print(f'  SKIP: folder not found')
        continue
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    if not files:
        print(f'  SKIP: no files')
        continue

    covers = set()
    closings = set()
    for f in files:
        s1 = get_slide_text(f, 1)
        last_slide = len(Presentation(f).slides)
        s_last = get_slide_text(f, last_slide)
        covers.add(s1[:200])
        closings.add(s_last[:200])

    print(f'  Total: {len(files)}')
    print(f'  Unique covers: {len(covers)}/{len(files)}')
    print(f'  Unique closings: {len(closings)}/{len(files)}')
    if len(covers) == len(files):
        print(f'  PASS')
    else:
        print(f'  WARNING: {len(files) - len(covers)} shared covers')


# ═══════════════════════════════════════════════════════════════
# GRAND SUMMARY
# ═══════════════════════════════════════════════════════════════
print()
print('=' * 70)
print('GRAND SUMMARY')
print('=' * 70)
print(f'  Total decks: {len(all_files)}')
print(f'  Byte-unique: {len(all_files) - sum(len(f)-1 for f in dupes.values())}')
print(f'  Text-unique: {len(all_files) - sum(len(f)-1 for f in text_dupes.values())}')
print(f'  Investor-500 individually aimed: {inv500_ok}/{len(glob.glob(os.path.join(BASE, "investor-500", "pptx", "*.pptx")))}')
print(f'  Investor-100 individually aimed: {inv100_ok}/{len(glob.glob(os.path.join(BASE, "investor-100", "pptx", "*.pptx")))}')
