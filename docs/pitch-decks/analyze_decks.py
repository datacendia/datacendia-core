"""Analyze all pitch deck PPTX files to understand slide structure."""
from pptx import Presentation
import os, glob

base = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'
folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

for folder in folders:
    path = os.path.join(base, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    for f in files[:3]:
        try:
            prs = Presentation(f)
            name = os.path.basename(f)
            print(f'\n=== {folder}/{name} --- {len(prs.slides)} slides ===')
            for i, slide in enumerate(prs.slides):
                texts = []
                has_image = False
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # picture
                        has_image = True
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t[:80])
                img_tag = ' [IMG]' if has_image else ''
                first_text = texts[0] if texts else '(no text)'
                print(f'  Slide {i+1}{img_tag}: {first_text}')
        except Exception as e:
            print(f'{folder}/{name} --- ERROR: {e}')
