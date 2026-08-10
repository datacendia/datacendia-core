"""Check every category for template bleed-through text."""
import os, sys, glob, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

def get_all_text(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except:
        return ''
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return '\n'.join(texts)


def check_bleed(folder, phrases, skip_filenames=None):
    path = os.path.join(BASE, folder, 'pptx')
    if not os.path.exists(path):
        return 0, 0, []
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    if skip_filenames:
        files = [f for f in files if not any(s in os.path.basename(f) for s in skip_filenames)]
    hits = []
    for f in files:
        text = get_all_text(f)
        for phrase in phrases:
            if phrase.lower() in text.lower():
                hits.append((os.path.basename(f), phrase))
                break
    return len(files), len(hits), hits


# 1. investor-500: a16z bleed
print('=== INVESTOR-500: a16z template bleed ===')
total, bad, hits = check_bleed('investor-500',
    ['dedicated AI fund', 'a16z', 'Andreessen Horowitz'])
print(f'  Files: {total}  Bleed: {bad}')
for fn, p in hits[:10]:
    print(f'    {fn}: "{p}"')
if len(hits) > 10:
    print(f'    ... +{len(hits)-10} more')

# 2. investor-100: a16z bleed (skip actual a16z deck)
print('\n=== INVESTOR-100: a16z template bleed ===')
total, bad, hits = check_bleed('investor-100',
    ['dedicated AI fund', 'a16z', 'Andreessen Horowitz'],
    skip_filenames=['a16z'])
print(f'  Files: {total}  Bleed: {bad}')
for fn, p in hits[:10]:
    print(f'    {fn}: "{p}"')
if len(hits) > 10:
    print(f'    ... +{len(hits)-10} more')

# 3. investor-50 extra (52+): Banking bleed
print('\n=== INVESTOR-50 (52+): Banking template bleed ===')
path = os.path.join(BASE, 'investor-50', 'pptx')
files = sorted(glob.glob(os.path.join(path, '*.pptx')))
extra = [f for f in files if os.path.basename(f)[:2].isdigit() and int(os.path.basename(f).split('-')[0]) >= 52]
banking_phrases = ['banking sector', 'BANKING', 'Basel III', 'Basel II', 'KYC/AML', 'bank compliance']
bad = 0
hits = []
for f in extra:
    text = get_all_text(f)
    for phrase in banking_phrases:
        if phrase.lower() in text.lower():
            bad += 1
            hits.append((os.path.basename(f), phrase))
            break
print(f'  Files: {len(extra)}  Bleed: {bad}')
for fn, p in hits[:10]:
    print(f'    {fn}: "{p}"')

# 4. sales: Fortune 500 bleed
print('\n=== SALES: Fortune 500 template bleed ===')
total, bad, hits = check_bleed('sales',
    ['Fortune 500', 'FORTUNE 500'],
    skip_filenames=['fortune'])
print(f'  Files: {total}  Bleed: {bad}')
for fn, p in hits[:5]:
    print(f'    {fn}: "{p}"')

# 5. design-partner: Healthcare bleed (skip healthcare deck)
print('\n=== DESIGN-PARTNER: Healthcare template bleed ===')
total, bad, hits = check_bleed('design-partner',
    ['HIPAA', 'clinical decision', 'patient safety', 'clinical AI'],
    skip_filenames=['healthcare'])
print(f'  Files: {total}  Bleed: {bad}')
for fn, p in hits[:10]:
    print(f'    {fn}: "{p}"')
if len(hits) > 10:
    print(f'    ... +{len(hits)-10} more')

# 6. enterprise: old banking/finance template bleed
print('\n=== ENTERPRISE: Banking template bleed ===')
total, bad, hits = check_bleed('enterprise',
    ['BANKING', 'Banking sector', 'Basel III', 'KYC/AML'])
print(f'  Files: {total}  Bleed: {bad}')

# 7. regional: old banking/finance template bleed
print('\n=== REGIONAL: Banking template bleed ===')
total, bad, hits = check_bleed('regional',
    ['BANKING', 'Banking sector', 'Basel III', 'KYC/AML'])
print(f'  Files: {total}  Bleed: {bad}')

# 8. Count how many investor-500 have the "dedicated AI fund" problem specifically
print('\n=== INVESTOR-500: "dedicated AI fund" specifically ===')
path = os.path.join(BASE, 'investor-500', 'pptx')
files = sorted(glob.glob(os.path.join(path, '*.pptx')))
bad = 0
for f in files:
    text = get_all_text(f)
    if 'dedicated ai fund' in text.lower():
        bad += 1
print(f'  Files: {len(files)}  With "dedicated AI fund": {bad}')

# 9. investor-100: "dedicated AI fund" specifically
print('\n=== INVESTOR-100: "dedicated AI fund" specifically ===')
path = os.path.join(BASE, 'investor-100', 'pptx')
files = sorted(glob.glob(os.path.join(path, '*.pptx')))
files = [f for f in files if 'a16z' not in os.path.basename(f)]
bad = 0
for f in files:
    text = get_all_text(f)
    if 'dedicated ai fund' in text.lower():
        bad += 1
print(f'  Files: {len(files)}  With "dedicated AI fund": {bad}')

print('\n=== SUMMARY ===')
print('Issues to fix are listed above.')
