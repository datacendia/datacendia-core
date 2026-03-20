"""
Generate personalized email text for each individual investor deck.
Extracts REAL data from actual PPTX files — no fabrication.
Outputs: docs/pitch-decks/INVESTOR-EMAILS.md
"""
import os, sys, io, glob, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE, 'INVESTOR-EMAILS.md')


def extract_investor_data(pptx_path):
    """Extract real firm data from a PPTX file."""
    try:
        prs = Presentation(pptx_path)
    except:
        return None

    data = {
        'filename': os.path.basename(pptx_path),
        'firm_name': '',
        'firm_type': '',
        'location': '',
        'aum': '',
        'portfolio': [],
        'thesis_angle': '',
        'alignment_text': '',
    }

    for si, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

        full = '\n'.join(texts)

        # Slide 1: Extract firm name, type, location, AUM
        if si == 0:
            for t in texts:
                if t.startswith('Datacendia for '):
                    data['firm_name'] = t.replace('Datacendia for ', '')
                # Look for the metadata line: "Type · Location · Stage · AUM"
                if '·' in t and ('AUM' in t or 'Capital' in t.split('·')[0] or 'Venture' in t or 'Equity' in t or 'Fund' in t or 'Bank' in t or 'Sovereign' in t or 'Accelerator' in t or 'Advisory' in t or 'Corporate' in t or 'Angel' in t or 'Pension' in t or 'Endowment' in t or 'Family' in t):
                    parts = [p.strip() for p in t.split('·')]
                    if len(parts) >= 2:
                        data['firm_type'] = parts[0]
                        data['location'] = parts[1] if len(parts) > 1 else ''
                        # Find AUM
                        for p in parts:
                            if '$' in p or 'AUM' in p or 'B+' in p or 'M+' in p:
                                data['aum'] = p.replace(' AUM', '').strip()

            # Fallback: try PREPARED FOR line
            if not data['firm_name']:
                for t in texts:
                    if 'PREPARED FOR' in t and t != 'DATACENDIA':
                        name = t.replace('PREPARED FOR ', '').strip()
                        if name and name != 'DATACENDIA':
                            data['firm_name'] = name.title()

        # Slide 2: Extract alignment text and portfolio
        if si == 1:
            for t in texts:
                if 'aligns with' in t.lower() or 'thesis' in t.lower() or 'mission' in t.lower():
                    if len(t) > 30:
                        data['alignment_text'] = t
                # Portfolio companies (short names, no long sentences)
                if len(t) < 30 and t not in ['DATACENDIA', 'INVESTOR ALIGNMENT', 'Portfolio Adjacency', 'Your Thesis → Our Mission'] and '©' not in t and 'datacendia' not in t.lower() and 'WHY' not in t:
                    if not any(x in t for x in ['has backed', 'completes', 'aligns']):
                        if t and t[0].isupper() and len(t) > 1:
                            data['portfolio'].append(t)

        # Slide 4: Extract thesis angle (Key Angle line)
        if si == 3:
            for t in texts:
                if 'Key Angle' in t:
                    continue
                if 'next infrastructure default' in t.lower() or 'decision evidence' in t.lower():
                    if len(t) > 20 and 'Sector' not in t:
                        data['thesis_angle'] = t

    # Clean up portfolio — remove duplicates and non-company items
    seen = set()
    clean_portfolio = []
    skip = {'DATACENDIA', 'INVESTOR ALIGNMENT', 'WHY', 'Portfolio Adjacency'}
    for p in data['portfolio']:
        if p not in seen and p not in skip and len(p) < 25:
            seen.add(p)
            clean_portfolio.append(p)
    data['portfolio'] = clean_portfolio[:6]  # Max 6

    return data


def fix_case(name):
    """Fix ALL CAPS names to proper case, preserving known acronyms."""
    acronyms = {'A16Z','GV','NEA','CRV','IVP','ADIA','GIC','KKR','DARPA','CIA','NSF','DOE','NASA','SPRIND','EIC','NGI','NTT','IBM','BECO','STV','GFC','JAFCO','DCM'}
    words = name.split()
    fixed = []
    for w in words:
        if w.upper() in acronyms:
            fixed.append(w.upper())
        elif w.isupper() and len(w) > 2:
            fixed.append(w.title())
        else:
            fixed.append(w)
    result = ' '.join(fixed)
    # Fix common patterns
    result = result.replace('(a16z)', '(a16z)').replace('(A16z)', '(a16z)')
    return result

def generate_email(data, category_folder):
    """Generate personalized email text from extracted data."""
    raw_firm = data['firm_name'] or data['filename'].replace('.pptx', '').split('-', 1)[-1].replace('-', ' ').title()
    firm = fix_case(raw_firm)
    firm_type = data['firm_type'] or 'Investor'
    location = data['location']
    aum = data['aum']
    portfolio = data['portfolio']
    thesis = data['thesis_angle']
    alignment = data['alignment_text']

    # Build the subject
    subject = f"Datacendia — Decision Evidence Infrastructure for AI Systems | Prepared for {firm}"

    # Build portfolio line
    portfolio_line = ''
    if portfolio:
        portfolio_line = f"Your portfolio includes {', '.join(portfolio[:3])}{',' if len(portfolio) > 3 else ''} — companies building the data and AI stack. Datacendia completes that stack with decision evidence and regulatory compliance for every AI decision."

    # Build location/type context
    context_line = ''
    parts = []
    if firm_type:
        parts.append(firm_type)
    if location:
        parts.append(location)
    if aum:
        parts.append(aum)
    if parts:
        context_line = f"*{' · '.join(parts)}*"

    # Build thesis alignment
    alignment_line = ''
    if thesis and 'decision evidence' in thesis.lower():
        alignment_line = thesis
    elif alignment and len(alignment) > 40:
        alignment_line = alignment

    # Build the email
    lines = []
    lines.append(f"### {firm}")
    lines.append(f"**Deck:** `{category_folder}/{data['filename']}`")
    if context_line:
        lines.append(context_line)
    lines.append('')
    lines.append(f"**Subject:** {subject}")
    lines.append('')
    lines.append('**Body:**')
    lines.append('')
    lines.append(f"Hi [FIRST NAME],")
    lines.append('')
    lines.append(f"I'm Stuart Rainey, founder of Datacendia. I've built the decision evidence layer for enterprise AI — the infrastructure that produces cryptographic proof of how AI-assisted decisions are made.")
    lines.append('')

    if alignment_line:
        lines.append(f"**Why {firm}:** {alignment_line}")
        lines.append('')

    if portfolio_line:
        lines.append(portfolio_line)
        lines.append('')

    lines.append("**What's built (not idea-stage):**")
    lines.append("- 205K+ passing tests, 156 validated API endpoints, live production platform")
    lines.append("- Multi-agent deliberation with cryptographic evidence (Ed25519, Merkle trees, RFC 3161)")
    lines.append("- Sovereign deployment: cloud, private cloud, on-prem, air-gapped")
    lines.append("- NVIDIA Inception member")
    lines.append("- Active design-partner conversations with a financial institution and Big 4 channel")
    lines.append('')
    lines.append("**The ask:** $1.5M pre-seed at $7M pre-money. Engineering risk is resolved — this is a GTM-stage investment. Funds go to first 3 design partners, SOC 2 Type II, and EU AI Act regulatory validation.")
    lines.append('')
    lines.append(f"I've attached a deck personalized for {firm}. Happy to schedule a 20-minute deep dive or live demo at your convenience.")
    lines.append('')
    lines.append("Best,")
    lines.append("Stuart Rainey")
    lines.append("Founder & CEO, Datacendia")
    lines.append("contact@datacendia.com | datacendia.com | app.datacendia.com")

    return '\n'.join(lines)


def main():
    all_emails = []

    # Process investor-100
    folder = os.path.join(BASE, 'investor-100', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    print(f'investor-100: {len(files)} files')
    for f in files:
        data = extract_investor_data(f)
        if data and data['firm_name']:
            email = generate_email(data, 'investor-100/pptx')
            all_emails.append(email)

    # Process investor-500
    folder = os.path.join(BASE, 'investor-500', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    print(f'investor-500: {len(files)} files')
    for f in files:
        data = extract_investor_data(f)
        if data and data['firm_name']:
            email = generate_email(data, 'investor-500/pptx')
            all_emails.append(email)

    # Also process investor (generic 15)
    folder = os.path.join(BASE, 'investor', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    print(f'investor: {len(files)} files')
    for f in files:
        data = extract_investor_data(f)
        if data and data['firm_name']:
            email = generate_email(data, 'investor/pptx')
            all_emails.append(email)

    print(f'\nTotal emails generated: {len(all_emails)}')

    # Write output
    with open(OUTPUT, 'w', encoding='utf-8') as out:
        out.write('# Personalized Investor Emails\n\n')
        out.write('Each email below is generated from the **actual content** of the corresponding pitch deck PPTX file.\n')
        out.write('Replace `[FIRST NAME]` with the actual contact name. Attach the **PDF version** of the deck.\n\n')
        out.write(f'**Total: {len(all_emails)} personalized emails**\n\n')
        out.write('---\n\n')

        for i, email in enumerate(all_emails):
            out.write(email)
            out.write('\n\n---\n\n')

    print(f'Written to {OUTPUT}')


if __name__ == '__main__':
    main()
