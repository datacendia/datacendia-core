"""
Embed platform screenshots into pitch deck PPTX files.

Rules:
  - ONLY replace existing image placeholders (never add images to text-only slides)
  - Exactly 1 screenshot per image slide, matched by slide content
  - Same position and size as the existing placeholder
  - Regenerate PDFs after (separate step)

Usage:
  python embed_screenshots.py --test          # Process ONE deck, open it for review
  python embed_screenshots.py --dry-run       # Show what would happen, change nothing
  python embed_screenshots.py                 # Process all 240 decks
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Emu
import os, glob, sys, subprocess

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'
SCREENSHOTS = os.path.join(BASE, 'screenshots')

# ─── Content-to-screenshot mapping (matched by slide text keywords) ──────────
# Each image slide has a clear content title. Map it to the RIGHT screenshot.
CONTENT_MAP = {
    'multi-agent decision engine': '02-council-deliberation.png',
    'decision audit trail': '05-decisions.png',
}


def get_slide_text_lower(slide):
    """Get all text from a slide as one lowercase string."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t.lower())
    return ' '.join(parts)


def match_screenshot(slide_text):
    """Return the correct screenshot filename for this slide's content, or None."""
    for keyword, screenshot in CONTENT_MAP.items():
        if keyword in slide_text:
            return screenshot
    return None


def process_deck(pptx_path, dry_run=False):
    """Replace image placeholders with correct screenshots in one PPTX file."""
    prs = Presentation(pptx_path)
    name = os.path.basename(pptx_path)
    changes = 0

    for i, slide in enumerate(prs.slides):
        # Only process slides that ALREADY have an image
        existing_images = [
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if not existing_images:
            continue

        slide_text = get_slide_text_lower(slide)
        screenshot_name = match_screenshot(slide_text)

        if not screenshot_name:
            continue

        img_path = os.path.join(SCREENSHOTS, screenshot_name)
        if not os.path.exists(img_path):
            print(f'  WARNING: {screenshot_name} not found in screenshots/')
            continue

        # Get the existing image's exact position and size
        old_img = existing_images[0]
        left = old_img.left
        top = old_img.top
        width = old_img.width
        height = old_img.height

        if dry_run:
            print(f'  Slide {i+1}: replace placeholder -> {screenshot_name} '
                  f'at ({round(left/914400,1)},{round(top/914400,1)}) '
                  f'{round(width/914400,1)}x{round(height/914400,1)}in')
        else:
            # Remove ALL existing images on this slide
            for img_shape in existing_images:
                sp = img_shape._element
                sp.getparent().remove(sp)

            # Add the correct screenshot at the exact same position/size
            slide.shapes.add_picture(img_path, left, top, width, height)

        changes += 1

    if changes > 0 and not dry_run:
        prs.save(pptx_path)

    return changes


def main():
    dry_run = '--dry-run' in sys.argv
    test_mode = '--test' in sys.argv

    if dry_run:
        print('DRY RUN -- no files will be modified\n')
    if test_mode:
        print('TEST MODE -- processing only 001-a16z.pptx\n')

    folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx',
               'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

    total_files = 0
    total_changes = 0

    for folder in folders:
        path = os.path.join(BASE, folder)
        if not os.path.exists(path):
            continue
        files = sorted(glob.glob(os.path.join(path, '*.pptx')))

        if test_mode:
            # Only process the a16z deck for testing
            files = [f for f in files if '001-a16z' in os.path.basename(f)]
            if not files:
                continue

        print(f'\n{folder}/ -- {len(files)} deck(s)')

        for f in files:
            name = os.path.basename(f)
            changes = process_deck(f, dry_run)
            if changes > 0:
                action = 'would replace' if dry_run else 'replaced'
                print(f'  {name}: {action} {changes} image(s)')
                total_changes += changes
            total_files += 1

    action = 'Would replace' if dry_run else 'Replaced'
    print(f'\nDONE: {total_files} files, {total_changes} images {action.lower()}')

    # In test mode, open the file for review
    if test_mode and not dry_run:
        test_file = os.path.join(BASE, 'investor-100', 'pptx', '001-a16z.pptx')
        if os.path.exists(test_file):
            print(f'\nOpening {test_file} for review...')
            subprocess.Popen(['start', '', test_file], shell=True)


if __name__ == '__main__':
    main()
