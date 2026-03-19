"""Get ALL text from each slide to understand content structure."""
from pptx import Presentation
from pptx.util import Inches, Emu
import os, glob

base = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

# Check one deck from each category in detail
samples = [
    'pptx/index.pptx',
    'investor/pptx/01-series-a-overview.pptx',
    'sales/pptx/01-fortune-500.pptx',
    'design-partner/pptx/01-healthcare.pptx',
    'gamma/pptx/01-investor-overview.pptx',
    'investor-100/pptx/001-a16z.pptx',
]

for rel in samples:
    f = os.path.join(base, rel)
    if not os.path.exists(f):
        continue
    prs = Presentation(f)
    print(f'\n{"="*80}')
    print(f'{rel} --- {len(prs.slides)} slides')
    print(f'{"="*80}')
    for i, slide in enumerate(prs.slides):
        print(f'\n--- Slide {i+1} ---')
        for shape in slide.shapes:
            stype = shape.shape_type
            if stype == 13:  # picture
                print(f'  [IMAGE] {shape.width}x{shape.height} at ({shape.left},{shape.top})')
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        print(f'  TEXT: {t[:120]}')

# Count total PPTX files
total = 0
for folder in ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']:
    path = os.path.join(base, folder)
    if os.path.exists(path):
        count = len(glob.glob(os.path.join(path, '*.pptx')))
        total += count
        print(f'\n{folder}: {count} files')
print(f'\nTOTAL PPTX FILES: {total}')
