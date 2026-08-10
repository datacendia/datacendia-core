"""Deep audit: extract full slide text from representative decks in every category."""
import os, sys, glob, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

def extract_all_text(pptx_path):
    """Return list of (slide_num, all_text) for every slide."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return [(0, f'ERROR reading: {e}')]
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        slides.append((i, '\n'.join(texts)))
    return slides

# Pick 2-3 representative files from each category
CATEGORIES = {
    'investor-100': 'investor-100/pptx',
    'investor-500': 'investor-500/pptx',
    'investor-50':  'investor-50/pptx',
    'investor':     'investor/pptx',
    'sales':        'sales/pptx',
    'design-partner': 'design-partner/pptx',
    'enterprise':   'enterprise/pptx',
    'regional':     'regional/pptx',
    'gamma':        'gamma/pptx',
    'main':         'pptx',
}

TEMPLATE_BLEED = [
    'a16z', 'Andreessen Horowitz', 'Fortune 500', 'FORTUNE 500',
    'Healthcare', 'HEALTHCARE', 'Banking', 'BANKING',
    'HIPAA', 'clinical', 'patient safety',
    'AI Governance for Regulated Financial Institutions',
    'DORA', '& Capital Markets',
]

for cat, folder in CATEGORIES.items():
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    if not files:
        continue
    # Sample: first, middle, last
    samples = [files[0]]
    if len(files) > 2:
        samples.append(files[len(files)//2])
    if len(files) > 1:
        samples.append(files[-1])
    
    print(f'\n{"="*80}')
    print(f'CATEGORY: {cat} ({len(files)} files)')
    print(f'{"="*80}')
    
    for fpath in samples:
        fname = os.path.basename(fpath)
        slides = extract_all_text(fpath)
        print(f'\n  --- {fname} ({len(slides)} slides) ---')
        for snum, text in slides:
            # Check for template bleed
            bleeds = []
            for token in TEMPLATE_BLEED:
                if token.lower() in text.lower():
                    # Only flag if it seems wrong for this category
                    if cat == 'investor-100' and token in ['a16z', 'Andreessen Horowitz']:
                        if 'a16z' not in fname:
                            bleeds.append(token)
                    elif cat == 'investor-500' and token in ['a16z', 'Andreessen Horowitz']:
                        bleeds.append(token)
                    elif cat == 'sales' and token in ['Fortune 500', 'FORTUNE 500']:
                        if 'fortune' not in fname:
                            bleeds.append(token)
                    elif cat == 'design-partner' and token in ['Healthcare', 'HEALTHCARE', 'HIPAA', 'clinical', 'patient safety']:
                        if 'healthcare' not in fname:
                            bleeds.append(token)
                    elif cat in ['enterprise', 'regional'] and token in ['Banking', 'BANKING', 'AI Governance for Regulated Financial Institutions', 'DORA', '& Capital Markets']:
                        bleeds.append(token)
                    elif cat == 'investor-50' and token in ['Banking', 'BANKING']:
                        if 'banking' not in fname:
                            bleeds.append(token)
            
            bleed_flag = f' *** BLEED: {bleeds}' if bleeds else ''
            # Truncate long text for readability
            preview = text[:300].replace('\n', ' | ')
            print(f'    Slide {snum}: {preview}{bleed_flag}')
