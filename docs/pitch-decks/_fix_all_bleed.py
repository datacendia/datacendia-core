"""
Fix ALL template bleed-through across every pitch deck category.

Issues addressed:
  1. investor-500 (380 decks): "dedicated AI fund" a16z bleed on Slide 2
  2. investor-50 extra (49 decks): Banking template bleed (Basel III, DORA, MiFID II, etc.)
  3. investor-100 (2 decks): minor a16z bleed
  4. enterprise (2 decks): minor banking bleed
"""
import contextlib
import io
import os
import sys
import glob
import re

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

# Suppress prints from imports
with contextlib.redirect_stdout(io.StringIO()):
    from _generate_1000 import (
        BASE, INVESTOR_500_BATCH2, INVESTOR_50_EXTRA,
        ENTERPRISE_TARGETS,
    )
    from _firms_database import INVESTOR_500_NEW

# Build lookup tables
ALL_INVESTOR_500 = INVESTOR_500_NEW + INVESTOR_500_BATCH2
INV500_BY_SLUG = {f['slug']: f for f in ALL_INVESTOR_500}
INV500_BY_NUM = {f['num']: f for f in ALL_INVESTOR_500}
INV50_BY_SLUG = {t['slug']: t for t in INVESTOR_50_EXTRA}
ENT_BY_SLUG = {t['slug']: t for t in ENTERPRISE_TARGETS}


# Also load batch 3 from _generate_remaining.py
try:
    with contextlib.redirect_stdout(io.StringIO()):
        from _generate_remaining import BATCH3
    for f in BATCH3:
        INV500_BY_SLUG[f['slug']] = f
        INV500_BY_NUM[f['num']] = f
except ImportError:
    pass


def replace_in_runs(shape, old_text, new_text, case_insensitive=False):
    """Replace text across runs in a shape. Returns True if any replacement made."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if case_insensitive:
                if old_text.lower() in run.text.lower():
                    run.text = re.sub(re.escape(old_text), new_text, run.text, flags=re.IGNORECASE)
                    changed = True
            else:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    changed = True
    return changed


def replace_paragraph_text(shape, old_fragment, new_full_text, case_insensitive=False):
    """Find paragraph containing old_fragment and replace ALL runs with new_full_text."""
    if not shape.has_text_frame:
        return False
    changed = False
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        match = old_fragment.lower() in full.lower() if case_insensitive else old_fragment in full
        if match:
            # Replace: put all text in first run, clear the rest
            if para.runs:
                para.runs[0].text = new_full_text
                for r in para.runs[1:]:
                    r.text = ''
                changed = True
    return changed


# ═══════════════════════════════════════════════════════════════════════════════
# 1. FIX INVESTOR-500: Remove "dedicated AI fund" bleed
# ═══════════════════════════════════════════════════════════════════════════════

def fix_investor_500():
    folder = os.path.join(BASE, 'investor-500', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    fixed = 0

    for fpath in files:
        fname = os.path.basename(fpath)
        # Extract slug from filename: "121-benchmark.pptx" -> "benchmark"
        parts = fname.replace('.pptx', '').split('-', 1)
        if len(parts) < 2:
            continue
        num_str, slug = parts
        try:
            num = int(num_str)
        except ValueError:
            continue

        firm = INV500_BY_NUM.get(num) or INV500_BY_SLUG.get(slug)
        if not firm:
            continue

        prs = Presentation(fpath)
        changed = False

        for slide in prs.slides:
            for shape in slide.shapes:
                # Fix 1: "dedicated AI fund" paragraph
                if replace_paragraph_text(
                    shape,
                    'dedicated AI fund',
                    f"{firm['name']} invests in {firm['type'].lower()} opportunities. Datacendia is the missing governance layer for enterprise AI.",
                    case_insensitive=True
                ):
                    changed = True

                # Fix 2: slug used as display name in alignment text
                # e.g., "benchmark's dedicated AI fund" -> already fixed above
                # But also fix "AI is eating software." prefix if still present
                if replace_paragraph_text(
                    shape,
                    'AI is eating software',
                    f"{firm['name']} invests in {firm['type'].lower()} opportunities. Datacendia is the missing governance layer for enterprise AI.",
                    case_insensitive=True
                ):
                    changed = True

        if changed:
            prs.save(fpath)
            fixed += 1

    print(f'  investor-500 fixed: {fixed}/{len(files)}')
    return fixed


# ═══════════════════════════════════════════════════════════════════════════════
# 2. FIX INVESTOR-50 EXTRA: Remove banking template bleed
# ═══════════════════════════════════════════════════════════════════════════════

# Generic replacements that work for any AI governance topic
BANKING_REPLACEMENTS = [
    # Slide 1 title bleed
    ('& Capital Markets', ''),
    ('AI Governance for Regulated Financial Institutions', 'Decision Evidence Infrastructure for AI Systems'),

    # Slide 2: regulatory references
    ('DORA', 'EU AI Act'),
    ('Jan 2025 — AI Operational Resilience', '2025–2026 — High-Risk AI Governance Mandated'),
    ('Basel III', 'ISO/IEC 42001'),
    ('Auditable Risk Models Required', 'AI Management System Certification'),
    ('MiFID II', 'NIST AI RMF'),
    ('Full Audit Trail for Algo Trading', 'Risk Management for Trustworthy AI'),
    ('AI Platforms w/ Crypto Decision Evidence', 'AI Platforms w/ Cryptographic Decision Evidence'),

    # Slide 3: problem statement
    ('Banks Deploy AI Without Audit Trails', 'Organizations Deploy AI Without Decision Evidence'),
    ('Credit decisions influenced by AI with no explainability for regulators',
     'High-stakes decisions influenced by AI with no explainability for stakeholders'),
    ('DORA requires operational resilience documentation for all AI systems by Jan 2025',
     'EU AI Act requires risk management and human oversight for high-risk AI systems'),
    ('Basel III mandates auditable risk models — AI black boxes fail this test',
     'ISO/IEC 42001 mandates AI management systems — black-box AI fails this test'),
    ('MiFID II requires full audit trails for algorithmic trading decisions',
     'Regulators and boards increasingly require full audit trails for AI-influenced decisions'),

    # Slide 4: solution
    ('credit, trading, and risk decision — regulator-ready proof before the auditor calls.',
     'high-stakes decision — regulator-ready proof before the auditor calls.'),
    ('Credit decision review with dissent tracking',
     'Decision review with dissent tracking'),
    ('M&A due diligence with full evidence trail',
     'Due diligence with full evidence trail'),
    ('DORA compliance with automated evidence',
     'Compliance with automated evidence generation'),
    ('Investment committee structured deliberation',
     'Committee-level structured deliberation'),
    ('Regulatory change mapping to P&L impact',
     'Regulatory change mapping to operational impact'),

    # Slide 5: thesis
    ('Banking is the largest regulated vertical. DORA enforcement creates immediate demand. Only platform providing cryptographic decision evidence — a regulatory requirement with zero incumbent solutions.',
     'AI governance is becoming mandatory across regulated industries. EU AI Act enforcement creates immediate demand. Only platform providing cryptographic decision evidence — a regulatory requirement with zero incumbent solutions.'),

    # Slide 6: founder
    ('Built the exact data infrastructure that financial institutions need for AI governance — at one of the world\'s largest asset managers.',
     'Built the data infrastructure and decision workflows required for regulated, high-stakes AI deployments at institutional scale.'),

    # Slide 10: closing
    # The title replacement already handled "Banking" -> theme name
]


def fix_investor_50_extra():
    folder = os.path.join(BASE, 'investor-50', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    fixed = 0

    for fpath in files:
        fname = os.path.basename(fpath)
        # Only fix files 52+ (the generated extras)
        num_str = fname.split('-')[0]
        try:
            num = int(num_str)
        except ValueError:
            continue
        if num < 52:
            continue

        prs = Presentation(fpath)
        changed = False

        for slide in prs.slides:
            for shape in slide.shapes:
                for old_text, new_text in BANKING_REPLACEMENTS:
                    if replace_in_runs(shape, old_text, new_text):
                        changed = True

        if changed:
            prs.save(fpath)
            fixed += 1

    print(f'  investor-50 extra fixed: {fixed}/49')
    return fixed


# ═══════════════════════════════════════════════════════════════════════════════
# 3. FIX INVESTOR-100: 2 decks with a16z bleed
# ═══════════════════════════════════════════════════════════════════════════════

def fix_investor_100():
    fixed = 0

    # 018-menlo.pptx
    fpath = os.path.join(BASE, 'investor-100', 'pptx', '018-menlo.pptx')
    if os.path.exists(fpath):
        prs = Presentation(fpath)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if replace_paragraph_text(shape, 'dedicated AI fund',
                    'Menlo Ventures invests in venture capital opportunities. Datacendia is the missing governance layer for enterprise AI.',
                    case_insensitive=True):
                    changed = True
                if replace_paragraph_text(shape, 'AI is eating software',
                    'Menlo Ventures invests in venture capital opportunities. Datacendia is the missing governance layer for enterprise AI.',
                    case_insensitive=True):
                    changed = True
        if changed:
            prs.save(fpath)
            fixed += 1

    # 091-andreessen-bio.pptx - "a16z" might appear legitimately if it's an a16z sub-fund
    fpath = os.path.join(BASE, 'investor-100', 'pptx', '091-andreessen-bio.pptx')
    if os.path.exists(fpath):
        prs = Presentation(fpath)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                # Only fix if slug "a16z" appears in non-name context
                if replace_in_runs(shape, "a16z's dedicated AI fund",
                    "a16z Bio invests in life-science and health-tech opportunities"):
                    changed = True
        if changed:
            prs.save(fpath)
            fixed += 1

    print(f'  investor-100 fixed: {fixed}/2')
    return fixed


# ═══════════════════════════════════════════════════════════════════════════════
# 4. FIX ENTERPRISE: 2 decks with banking bleed
# ═══════════════════════════════════════════════════════════════════════════════

def fix_enterprise():
    folder = os.path.join(BASE, 'enterprise', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    fixed = 0

    for fpath in files:
        prs = Presentation(fpath)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += ' '.join(p.text for p in shape.text_frame.paragraphs)

        if 'BANKING' not in text and 'Banking sector' not in text and 'Basel III' not in text and 'KYC/AML' not in text:
            continue

        prs = Presentation(fpath)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                for old_text, new_text in BANKING_REPLACEMENTS:
                    if replace_in_runs(shape, old_text, new_text):
                        changed = True
        if changed:
            prs.save(fpath)
            fixed += 1

    print(f'  enterprise fixed: {fixed}')
    return fixed


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    print('Fixing all template bleed-through...\n')

    total = 0
    total += fix_investor_500()
    total += fix_investor_50_extra()
    total += fix_investor_100()
    total += fix_enterprise()

    print(f'\nTotal decks fixed: {total}')
    print('Done.')


if __name__ == '__main__':
    main()
