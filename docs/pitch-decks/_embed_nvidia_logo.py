"""Embed the NVIDIA Inception badge on every slide that mentions 'NVIDIA Inception Member'."""
import os, sys, io, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Emu

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'
NVIDIA_BADGE = r'C:\Users\Stu\datacendia-core\docs\Inception Badges\for-screen\nvidia-inception-program-badge-rgb-for-screen.png'

CATEGORIES = [
    'investor-500/pptx', 'investor-100/pptx', 'investor-50/pptx',
    'investor/pptx', 'sales/pptx', 'design-partner/pptx',
    'enterprise/pptx', 'regional/pptx', 'gamma/pptx', 'pptx',
]


def embed_nvidia_badge(pptx_path):
    """Find slides with 'NVIDIA Inception Member' text and add the badge image."""
    prs = Presentation(pptx_path)
    changed = False

    for slide in prs.slides:
        # Check if this slide has "NVIDIA Inception Member" text
        has_nvidia_text = False
        nvidia_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if 'NVIDIA Inception Member' in para.text:
                        has_nvidia_text = True
                        nvidia_shape = shape
                        break
            if has_nvidia_text:
                break

        if not has_nvidia_text or nvidia_shape is None:
            continue

        # Check if a badge image already exists on this slide
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        existing_images = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        # Skip if there are already images that might be the badge
        # (screenshots are large, badge would be small)
        badge_exists = False
        for img in existing_images:
            if img.width < Inches(2.5) and img.height < Inches(1.5):
                badge_exists = True
                break

        if badge_exists:
            continue

        # Position the badge to the left of the "NVIDIA Inception Member" text
        # The text shape gives us a reference point
        text_left = nvidia_shape.left
        text_top = nvidia_shape.top

        # Place badge to the left of the text, vertically centered
        badge_width = Inches(1.8)
        badge_height = Inches(0.9)
        badge_left = text_left - badge_width - Inches(0.15)
        badge_top = text_top + (nvidia_shape.height - badge_height) // 2

        # Ensure badge doesn't go off-screen
        if badge_left < Inches(0.2):
            badge_left = Inches(0.2)

        slide.shapes.add_picture(NVIDIA_BADGE, badge_left, badge_top, badge_width, badge_height)
        changed = True

    if changed:
        prs.save(pptx_path)
    return changed


def main():
    if not os.path.exists(NVIDIA_BADGE):
        print(f'ERROR: NVIDIA badge not found at {NVIDIA_BADGE}')
        return

    total = 0
    for cat_folder in CATEGORIES:
        folder = os.path.join(BASE, cat_folder)
        if not os.path.exists(folder):
            continue
        files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
        cat_fixed = 0
        for f in files:
            if embed_nvidia_badge(f):
                cat_fixed += 1
        if cat_fixed > 0:
            print(f'  {cat_folder}: {cat_fixed}/{len(files)} badges added')
        total += cat_fixed

    print(f'\nTotal decks with NVIDIA badge added: {total}')


if __name__ == '__main__':
    main()
