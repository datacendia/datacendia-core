"""
Fix all pitch decks based on investor feedback.
Applies to ALL 1,000 PPTX files across all categories.

Changes:
1. "0 Platforms w/ crypto AI evidence" → defensible statement
2. "AI governance infrastructure" → "Decision evidence infrastructure"
3. "90-Day Guarantee" block → removed/replaced with "Why Now" inevitability
4. "2 warm enterprise intros" → "Active design-partner conversations"
5. Add moat language: "We define the artifact regulators will eventually expect"
6. Tighten the ask with outcome-based milestones
7. Sharpen positioning throughout
"""
import os, sys, io, glob, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

BASE = os.path.dirname(os.path.abspath(__file__))

# ─── Text replacements (old → new) applied to ALL shapes in ALL slides ────────
REPLACEMENTS = [
    # 1. Fix overclaiming — the most damaging line
    (
        'No existing platform provides cryptographic proof of AI-assisted decisions',
        'No widely adopted platform provides cryptographically verifiable decision evidence at the decision level'
    ),
    (
        '0\nPlatforms w/ crypto AI evidence',
        '0\nWidely adopted decision-evidence platforms'
    ),
    # Handle the stat box version
    (
        'Platforms w/ crypto AI evidence',
        'Widely adopted decision-evidence platforms'
    ),

    # 2. Sharpen positioning — governance → decision evidence
    (
        'Decision Crisis Immunization Infrastructure',
        'Decision Evidence Infrastructure for AI Systems'
    ),
    (
        'the missing governance layer for enterprise AI',
        'the missing decision-evidence layer for enterprise AI'
    ),
    (
        'AI governance is the next infrastructure default',
        'Decision evidence is the next infrastructure default — every enterprise deploying AI will be required to produce decision evidence under regulatory or legal challenge'
    ),
    (
        'governance, accountability, and regulatory\ncompliance for every AI decision',
        'decision evidence, accountability, and regulatory compliance for every AI decision'
    ),
    (
        'governance, accountability, and regulatory compliance for every AI decision',
        'decision evidence, accountability, and regulatory compliance for every AI decision'
    ),

    # 3. Reframe traction
    (
        '2 warm enterprise intros (financial institution, Big 4 channel). NVIDIA Inception member.',
        'Active enterprise design-partner conversations (financial institution, Big 4 channel). NVIDIA Inception member. 3 LOI-stage discussions.'
    ),
    (
        '2 warm enterprise intros (financial institution, Big 4 channel)',
        'Active enterprise design-partner conversations (financial institution, Big 4 channel)'
    ),
    (
        '2 warm enterprise intros',
        'Active enterprise design-partner conversations'
    ),

    # 4. Tighten the ask — add outcome milestones
    (
        'GTM hiring, SOC 2 Type II certification, first enterprise pilots.',
        'First 3 design partners signed, SOC 2 Type II certification, regulatory validation (EU AI Act compliance audit).'
    ),
    (
        'GTM hiring, SOC 2 Type II, first enterprise pilots',
        'First 3 design partners, SOC 2 Type II, regulatory validation'
    ),

    # 5. Remove/replace 90-day guarantee with "Why Now" inevitability
    (
        '90-Day Guarantee',
        'Why This Must Exist Now'
    ),
    (
        'Deploy Datacendia for 90 days on a single strategic problem. If the Council doesn\'t identify at least 3 critical risks you missed, we refund the pilot fee.',
        'EU AI Act (2025–2026) → enforceable auditability requirements. DORA → operational traceability for financial services. Litigation risk → retrospective scrutiny of AI decisions. Enterprise AI adoption → exploding unaudited decision surface. This creates a new infrastructure requirement: decision evidence.'
    ),
    # Handle truncated version
    (
        'Deploy Datacendia for 90 days on a single strategic problem. If the Council doesn\'t identify at least 3 critical risks',
        'EU AI Act → enforceable auditability. DORA → operational traceability. Litigation risk → retrospective scrutiny. This creates a new infrastructure requirement: decision evidence'
    ),

    # 6. Add moat language where "Key Angle" appears
    (
        'Sector: AI Infrastructure',
        'Sector: Decision Evidence Infrastructure\n\nOur moat: We define the evidentiary artifact that regulators will eventually expect. First to codify it wins the category.'
    ),
    (
        'Sector: AI Governance',
        'Sector: Decision Evidence Infrastructure\n\nOur moat: We define the evidentiary artifact that regulators will eventually expect.'
    ),

    # 7. Sharpen "Enterprise AI Has No Accountability Layer"
    (
        'Enterprise AI Has No Accountability Layer',
        'Enterprise AI Has No Decision Evidence Layer'
    ),
    (
        'AI systems make high-stakes decisions with no audit trail, no evidence, and no dissent tracking',
        'Every enterprise deploying AI will be required to produce decision evidence under regulatory or legal challenge — and no platform exists to generate it'
    ),

    # 8. Fix the "Let's Talk" slide positioning
    (
        'Personalized for',
        'Prepared for'
    ),

    # 9. Clean up any remaining "AI governance" to "decision evidence" in key positions
    (
        'AI governance infrastructure',
        'decision evidence infrastructure'
    ),
    (
        'AI Governance Infrastructure',
        'Decision Evidence Infrastructure'
    ),
]


def apply_replacements(prs_path):
    """Apply all text replacements to a single PPTX file."""
    try:
        prs = Presentation(prs_path)
    except Exception as e:
        print(f'  ERROR opening {prs_path}: {e}')
        return False

    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original = run.text
                    for old, new in REPLACEMENTS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                    if run.text != original:
                        changed = True

            # Also check full paragraph text (some text spans multiple runs)
            full_text = shape.text_frame.text
            for old, new in REPLACEMENTS:
                if old in full_text:
                    # Try to replace in individual runs
                    for para in shape.text_frame.paragraphs:
                        combined = ''.join(r.text for r in para.runs)
                        if old in combined:
                            new_combined = combined.replace(old, new)
                            if para.runs:
                                # Put all text in first run, clear the rest
                                font = para.runs[0].font
                                para.runs[0].text = new_combined
                                for r in para.runs[1:]:
                                    r.text = ''
                                changed = True

    if changed:
        prs.save(prs_path)
    return changed


def main():
    # Find all PPTX files
    folders = [
        'pptx', 'investor/pptx', 'investor-50/pptx', 'investor-100/pptx',
        'investor-500/pptx', 'sales/pptx', 'design-partner/pptx',
        'gamma/pptx', 'enterprise/pptx', 'regional/pptx'
    ]

    total = 0
    modified = 0

    for folder in folders:
        path = os.path.join(BASE, folder)
        if not os.path.exists(path):
            continue
        pptx_files = sorted(glob.glob(os.path.join(path, '*.pptx')))
        if not pptx_files:
            continue

        print(f'\n{folder} — {len(pptx_files)} files')
        for f in pptx_files:
            total += 1
            was_changed = apply_replacements(f)
            if was_changed:
                modified += 1
            if total % 100 == 0:
                print(f'  processed {total}...')

    print(f'\n{"="*60}')
    print(f'DONE: {total} files processed, {modified} modified')
    print(f'{"="*60}')


if __name__ == '__main__':
    main()
