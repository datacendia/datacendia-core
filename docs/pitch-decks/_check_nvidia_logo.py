"""Check how NVIDIA Inception is referenced across pitch decks."""
import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

# Check one deck from each major category
samples = [
    'investor-100/pptx/001-a16z.pptx',
    'investor-500/pptx/121-benchmark.pptx',
    'design-partner/pptx/01-healthcare.pptx',
    'investor-50/pptx/01-banking.pptx',
    'sales/pptx/01-fortune-500.pptx',
    'enterprise/pptx/01-ai-risk-management.pptx',
    'regional/pptx/01-united-states-northeast.pptx',
]

for rel in samples:
    fpath = os.path.join(BASE, rel)
    if not os.path.exists(fpath):
        continue
    prs = Presentation(fpath)
    print(f'\n=== {rel} ({len(prs.slides)} slides) ===')
    for i, slide in enumerate(prs.slides):
        slide_text = ' '.join(p.text for sh in slide.shapes if sh.has_text_frame for p in sh.text_frame.paragraphs)
        if 'nvidia' in slide_text.lower() or 'NVIDIA' in slide_text:
            print(f'  Slide {i+1}: NVIDIA text found')
            # List all shapes on this slide
            for sh in slide.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f'    IMAGE: {sh.name} pos=({sh.left},{sh.top}) size=({sh.width},{sh.height})')
                    # Check if image has content
                    try:
                        img = sh.image
                        print(f'      content_type={img.content_type} size={len(img.blob)} bytes')
                    except Exception as e:
                        print(f'      ERROR reading image: {e}')
                elif sh.has_text_frame:
                    txt = ' '.join(p.text for p in sh.text_frame.paragraphs).strip()
                    if txt:
                        print(f'    TEXT: {sh.name} = {txt[:120]}')
                elif sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    print(f'    PLACEHOLDER: {sh.name}')
                else:
                    print(f'    OTHER: {sh.name} type={sh.shape_type}')
