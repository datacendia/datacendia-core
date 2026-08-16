"""
Extract target investor/company info from all pitch decks.
Builds a contact directory from the deck content itself.
"""
from pptx import Presentation
import os, glob, json, re

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

def extract_deck_info(pptx_path):
    """Extract target company info from a deck."""
    prs = Presentation(pptx_path)
    rel = os.path.relpath(pptx_path, BASE)
    
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        all_text.append(t)
    
    combined = '\n'.join(all_text)
    
    info = {
        'file': rel,
        'filename': os.path.basename(pptx_path),
        'category': rel.split('\\')[0] if '\\' in rel else rel.split('/')[0],
        'all_text_sample': combined[:3000],
    }
    
    # Extract "PREPARED FOR" line
    for t in all_text:
        if 'prepared for' in t.lower():
            info['prepared_for'] = t
            break
    
    # Extract investor/company details from slide 1 and 2
    for t in all_text:
        tl = t.lower()
        # Look for sector/location info
        if any(kw in tl for kw in ['venture capital', 'private equity', 'sovereign wealth', 
               'family office', 'corporate venture', 'growth equity', 'seed to',
               'series a', 'early stage', 'late stage', 'impact invest',
               'deep tech', 'cybersecurity', 'fintech', 'regtech', 'ai fund',
               'defense', 'government', 'foundation', 'non-profit']):
            if 'aum' in tl or '$' in t or 'capital' in tl or 'venture' in tl or 'fund' in tl:
                info.setdefault('investor_description', t)
    
    # Look for "Why X" or "Why Datacendia Aligns with X" 
    for t in all_text:
        if 'why datacendia aligns with' in t.lower():
            info['alignment_reason'] = t
            break
    
    # Look for portfolio companies
    for t in all_text:
        if 'portfolio adjacency' in t.lower():
            # Next few lines are portfolio companies
            idx = all_text.index(t)
            companies = []
            for j in range(idx+1, min(idx+10, len(all_text))):
                c = all_text[j].strip()
                if c and len(c) < 50 and c not in ('DATACENDIA',) and '©' not in c:
                    companies.append(c)
                else:
                    break
            info['portfolio_companies'] = companies
            break
    
    # Look for thesis/angle
    for t in all_text:
        tl = t.lower()
        if 'your thesis' in tl or 'key angle' in tl:
            idx = all_text.index(t)
            if idx + 1 < len(all_text):
                info['thesis_angle'] = all_text[idx + 1] if all_text[idx+1] != 'DATACENDIA' else t
            break
    
    return info

# Process all decks
folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 
           'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

all_info = []
for folder in folders:
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    for f in files:
        try:
            info = extract_deck_info(f)
            all_info.append(info)
        except Exception as e:
            print(f'ERROR: {f}: {e}')

# Save full results
output = os.path.join(BASE, '_contact_directory.json')
with open(output, 'w', encoding='utf-8') as fp:
    json.dump(all_info, fp, indent=2, ensure_ascii=False)

# Print summary for investor-100 (most detailed)
print(f'Extracted info from {len(all_info)} decks\n')

for info in all_info:
    if 'investor-100' in info['file'] or 'investor-50' in info['file']:
        prepared = info.get('prepared_for', '(no PREPARED FOR line)')
        desc = info.get('investor_description', '')
        portfolio = ', '.join(info.get('portfolio_companies', [])[:5])
        print(f'{info["filename"]}')
        print(f'  Target: {prepared}')
        if desc:
            print(f'  Profile: {desc[:120]}')
        if portfolio:
            print(f'  Portfolio: {portfolio}')
        print()
