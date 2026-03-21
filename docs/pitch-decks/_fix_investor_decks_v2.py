"""
Comprehensive fix for ALL investor-500 (and investor-100) decks.
Works at the individual run level to handle text split across XML runs.

Fixes:
  1. "dedicated AI fund" a16z bleed (cross-run)
  2. a16z portfolio company names bleeding into other decks
  3. Fabricated adjacency text ("has backed the data stack")
  4. "AI is eating software" template phrase
  5. Slug used as display name (e.g. "bg-fund's" instead of firm name)
"""
import contextlib
import io
import os
import sys
import glob
import re

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

with contextlib.redirect_stdout(io.StringIO()):
    from _generate_1000 import BASE, INVESTOR_500_BATCH2
    from _firms_database import INVESTOR_500_NEW

try:
    with contextlib.redirect_stdout(io.StringIO()):
        from _generate_remaining import BATCH3
except Exception:
    BATCH3 = []

ALL_FIRMS = INVESTOR_500_NEW + INVESTOR_500_BATCH2 + BATCH3
INV_BY_NUM = {}
INV_BY_SLUG = {}
for f in ALL_FIRMS:
    INV_BY_NUM[f['num']] = f
    INV_BY_SLUG[f['slug']] = f

A16Z_PORTFOLIO = ['Databricks', 'Anyscale', 'Hex', 'dbt Labs', 'Replit']

VAGUE_WORDS = [
    'enterprise', 'startups', 'tech', 'ecosystem', 'investments',
    'portfolio', 'alumni', 'infrastructure', 'network', 'digital',
    'sector', 'region', 'market', 'industry', 'investing', 'lending',
    'opportunities', 'focused', 'opportunistic', 'legacy', 'allocation',
    'initiatives', 'economy', 'companies', 'globally',
]


def is_vague(entry):
    el = entry.lower()
    return any(v in el for v in VAGUE_WORDS) or len(entry) > 50


def real_portfolio(firm):
    return [p for p in firm.get('portfolio', []) if not is_vague(p)]


def fix_slide2_paragraph(para_text, firm):
    """Fix the full concatenated paragraph text for Slide 2 issues."""
    name = firm['name']
    slug = firm['slug']
    ftype = firm['type'].lower()
    rp = real_portfolio(firm)
    thesis = firm.get('thesis', 'AI governance infrastructure')
    changed = False
    original = para_text

    # Fix "dedicated AI fund"
    if 'dedicated AI fund' in para_text or 'dedicated ai fund' in para_text.lower():
        para_text = re.sub(
            r".*dedicated AI fund[^.]*\.?",
            name + ' invests in ' + ftype + ' opportunities. Datacendia is the missing governance layer for enterprise AI.',
            para_text, flags=re.IGNORECASE
        )
        changed = True

    # Fix "AI is eating software"
    if 'AI is eating software' in para_text:
        para_text = re.sub(r"AI is eating software\.\s*", "", para_text)
        changed = True

    # Fix slug-as-name (e.g., "bg-fund's")
    slug_poss = slug + "'s"
    if slug_poss in para_text.lower():
        para_text = re.sub(re.escape(slug) + r"'s", name + "'s", para_text, flags=re.IGNORECASE)
        changed = True

    # Fix adjacency text: "X has backed the data stack (...) and AI infra (...)"
    adj_match = re.search(r"has backed the data stack \([^)]*\) and AI infra \([^)]*\)\.[^.]*\.", para_text)
    if adj_match:
        if rp:
            pstr = ', '.join(rp[:4])
            repl = "portfolio includes " + pstr + ". Datacendia completes the stack with AI governance and audit."
        else:
            repl = "investment focus aligns with Datacendia. AI governance completes the enterprise stack."
        para_text = para_text[:adj_match.start()] + repl + para_text[adj_match.end():]
        changed = True

    # Fix a16z portfolio company names used as bullet points
    for old_co in A16Z_PORTFOLIO:
        if old_co in para_text:
            if rp:
                new_co = rp[A16Z_PORTFOLIO.index(old_co) % len(rp)] if A16Z_PORTFOLIO.index(old_co) < len(rp) else ''
            else:
                new_co = ''
            if new_co:
                para_text = para_text.replace(old_co, new_co)
            else:
                para_text = para_text.replace(old_co, '')
            changed = True

    return para_text, changed


def rewrite_paragraph_runs(para, new_text):
    """Put all text into first run, clear the rest."""
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ''


def fix_deck(fpath, firm):
    """Fix a single investor deck. Returns True if any changes made."""
    prs = Presentation(fpath)
    changed = False
    name = firm['name']
    slug = firm['slug']
    rp = real_portfolio(firm)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in para.runs)
                if not full.strip():
                    continue

                needs_fix = False
                # Check if this paragraph has any issues
                if 'dedicated AI fund' in full.lower():
                    needs_fix = True
                if 'AI is eating software' in full:
                    needs_fix = True
                if slug + "'s" in full.lower():
                    needs_fix = True
                if 'has backed the data stack' in full:
                    needs_fix = True
                # Check for a16z portfolio companies in non-a16z decks
                for co in A16Z_PORTFOLIO:
                    if co in full:
                        needs_fix = True
                        break

                if needs_fix:
                    new_text, was_changed = fix_slide2_paragraph(full, firm)
                    if was_changed:
                        rewrite_paragraph_runs(para, new_text)
                        changed = True

    if changed:
        prs.save(fpath)
    return changed


def main():
    print('Comprehensive investor deck fix (v2)...\n')

    # Fix investor-500
    folder = os.path.join(BASE, 'investor-500', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    fixed = 0
    for fpath in files:
        fname = os.path.basename(fpath)
        parts = fname.replace('.pptx', '').split('-', 1)
        if len(parts) < 2:
            continue
        try:
            num = int(parts[0])
        except ValueError:
            continue
        slug = parts[1]
        firm = INV_BY_NUM.get(num) or INV_BY_SLUG.get(slug)
        if not firm:
            continue
        if fix_deck(fpath, firm):
            fixed += 1
    print(f'  investor-500: {fixed}/{len(files)} fixed')

    # Fix investor-100 (skip a16z itself)
    folder = os.path.join(BASE, 'investor-100', 'pptx')
    files = sorted(glob.glob(os.path.join(folder, '*.pptx')))
    fixed100 = 0
    for fpath in files:
        fname = os.path.basename(fpath)
        if 'a16z' in fname and 'crypto' not in fname and 'growth' not in fname and 'bio' not in fname:
            continue
        # For investor-100, we need firm data too
        parts = fname.replace('.pptx', '').split('-', 1)
        if len(parts) < 2:
            continue
        try:
            num = int(parts[0])
        except ValueError:
            continue
        slug = parts[1]
        # investor-100 firms are numbered 001-120, not in our 121+ database
        # But we still fix a16z template bleed generically
        prs = Presentation(fpath)
        ch = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    full = ''.join(r.text for r in para.runs)
                    if 'dedicated AI fund' in full.lower():
                        new = re.sub(
                            r".*dedicated AI fund[^.]*\.?",
                            "This firm invests in technology opportunities. Datacendia is the missing governance layer for enterprise AI.",
                            full, flags=re.IGNORECASE
                        )
                        rewrite_paragraph_runs(para, new)
                        ch = True
                    if 'AI is eating software' in full:
                        new = full.replace('AI is eating software. ', '').replace('AI is eating software.', '')
                        rewrite_paragraph_runs(para, new)
                        ch = True
        if ch:
            prs.save(fpath)
            fixed100 += 1
    print(f'  investor-100: {fixed100}/{len(files)} fixed')

    print('\nDone.')


if __name__ == '__main__':
    main()
