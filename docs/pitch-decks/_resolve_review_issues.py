"""Resolve the issues found in the full pitch deck review."""
import contextlib
import io
import os
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

with contextlib.redirect_stdout(io.StringIO()):
    from _generate_1000 import BASE, SCREENSHOTS, DESIGN_PARTNER_EXTRA, ENTERPRISE_TARGETS, REGIONAL_TARGETS


REGIONAL_EXTRA = [
    {'num': 121, 'slug': 'east-african-community', 'name': 'East African Community - Arusha'},
    {'num': 122, 'slug': 'sadc-gaborone', 'name': 'SADC - Gaborone'},
]

DESIGN_PARTNER_BY_SLUG = {item['slug']: item for item in DESIGN_PARTNER_EXTRA}
ENTERPRISE_BY_SLUG = {item['slug']: item for item in ENTERPRISE_TARGETS}
REGIONAL_BY_SLUG = {item['slug']: item for item in REGIONAL_TARGETS + REGIONAL_EXTRA}

MULTILATERAL_TOKENS = [
    'EU - ', 'NATO', 'UN - ', 'World Bank', 'IMF', 'OECD', 'WEF', 'ASEAN',
    'African Union', 'GCC', 'APEC', 'G7', 'G20', 'BRICS', 'Pacific Islands Forum',
    'Caribbean Community', 'Nordic Council', 'Mercosur', 'SCO', 'Arab League',
]

EUROPE_TOKENS = ['Germany', 'France', 'Switzerland', 'Netherlands', 'Belgium', 'Luxembourg', 'Sweden', 'Norway', 'Denmark', 'Finland', 'Ireland', 'Spain', 'Portugal', 'Italy', 'Austria', 'Poland', 'Czech Republic', 'Estonia', 'Latvia', 'Lithuania', 'Romania', 'Greece', 'United Kingdom']
ASIA_TOKENS = ['Japan', 'South Korea', 'China', 'Hong Kong', 'Taiwan', 'Singapore', 'India', 'Indonesia', 'Malaysia', 'Thailand', 'Vietnam', 'Philippines']
AFRICA_TOKENS = ['South Africa', 'Nigeria', 'Kenya', 'Ghana', 'Rwanda', 'Morocco', 'Ethiopia', 'Tanzania']
LATAM_TOKENS = ['Brazil', 'Mexico', 'Colombia', 'Argentina', 'Chile', 'Peru', 'Costa Rica', 'Panama', 'Uruguay', 'Ecuador']
MIDDLE_EAST_TOKENS = ['Israel', 'Turkey', 'UAE', 'Saudi Arabia', 'Qatar', 'Bahrain', 'Kuwait', 'Oman', 'Jordan', 'Egypt']
NORTH_AMERICA_TOKENS = ['United States', 'Canada']

GAMMA_CONTENT = {
    '07-healthcare-vertical.pptx': {
        'cover_label': 'HEALTHCARE VERTICAL',
        'headline': 'Decision evidence for clinical, operational, and compliance AI',
        'subheadline': 'Hospitals, payers, and life-science teams need replayable reasoning before AI touches care pathways, utilization, or regulated workflows.',
        'problem_title': 'Healthcare AI Cannot Be a Black Box.',
        'problem_lines': [
            'Clinical, operational, and revenue-cycle decisions increasingly involve AI.',
            'Providers need evidence for triage, staffing, utilization, and exception handling.',
            'Patient safety, human override, and auditability are non-negotiable.',
            'Very few platforms preserve dissent, approvals, and replay in one system.',
            'Healthcare organizations are moving AI into high-stakes workflows faster than governance is catching up.',
            'Single-model outputs are not enough when teams need challenge, escalation, and accountable approvals.',
            'Datacendia adds a defensible record of how AI-informed decisions were made.',
        ],
        'market_title': 'Why Healthcare Governance Matters Now',
        'market_lines': [
            'Healthcare buyers want accountable AI before clinical, operational, or billing risk compounds.',
            'Auditability matters across hospital operations, payer workflows, and regulated life-science environments.',
            'Governance spend follows trust, incident preparedness, and explainable oversight.',
            'Decision evidence becomes the wedge into one of the world\'s largest regulated markets.',
            'Multi-agent review reduces blind spots before AI affects patients, reimbursements, or compliance posture.',
            'Deployment flexibility matters across cloud, on-prem, and restricted environments.',
            'Evidence export shortens legal, compliance, and quality-review cycles.',
            'Healthcare AI adoption will favor systems that can be defended after the fact.',
        ],
        'compliance_title': 'Healthcare Readiness',
        'compliance_body': 'Why this matters for fundraising: healthcare buyers care less about model novelty than about accountable deployment. Decision evidence supports internal review, compliance sign-off, incident replay, and trust across clinical and operational stakeholders.',
    },
    '08-financial-compliance.pptx': {
        'cover_label': 'FINANCIAL COMPLIANCE',
        'headline': 'Decision evidence for regulated finance and operational resilience',
        'subheadline': 'Banks, insurers, and capital-markets teams need replayable reasoning for credit, fraud, trading, and compliance decisions before scrutiny arrives.',
        'problem_title': 'Regulated Finance Needs Defensible AI Decisions.',
        'problem_lines': [
            'Financial institutions are moving AI into underwriting, surveillance, fraud, and operations.',
            'Supervisors and internal risk teams need evidence for who approved what and why.',
            'Operational resilience depends on replayable decisions, not just logs and dashboards.',
            'Very few platforms combine deliberation, dissent, and immutable evidence.',
            'AI adoption is accelerating across regulated workflows while explainability expectations keep tightening.',
            'Single-model outputs are too thin for risk committees, compliance teams, and external review.',
            'Datacendia closes the gap between AI output and a defensible decision record.',
        ],
        'market_title': 'Why Financial Governance Matters Now',
        'market_lines': [
            'Regulated finance rewards platforms that shorten review cycles without weakening controls.',
            'Operational resilience, model risk, and audit expectations are all pushing toward evidence-rich workflows.',
            'Credit, fraud, and surveillance teams need more than predictions — they need explainable approvals.',
            'Decision evidence creates a system of record for high-stakes financial workflows.',
            'Multi-agent challenge helps surface blind spots before they become losses or enforcement issues.',
            'Sovereign deployment matters for sensitive data and institution-specific controls.',
            'Exportable evidence improves regulator response, board reporting, and post-incident analysis.',
            'Finance AI winners will be the platforms institutions can defend under scrutiny.',
        ],
        'compliance_title': 'Financial Readiness',
        'compliance_body': 'Why this matters for fundraising: regulated finance buys trust, control, and auditability. Decision evidence turns AI from a model-risk concern into a governed operating capability across compliance, risk, and front-office workflows.',
    },
    '09-ai-governance-market.pptx': {
        'cover_label': 'AI GOVERNANCE MARKET',
        'headline': 'The category view: decision evidence is the missing AI layer',
        'subheadline': 'Model access is commoditizing. The durable category is the system that preserves reasoning, approvals, dissent, and replay for every material AI decision.',
        'problem_title': 'The Market Has Models, But Not Accountability.',
        'problem_lines': [
            'Enterprises can buy models, copilots, and orchestration tools almost anywhere.',
            'What they still lack is a defensible record of how high-stakes AI decisions were made.',
            'Auditability, replay, and preserved dissent remain missing from most deployments.',
            'That leaves boards, regulators, and operators exposed the moment a decision is challenged.',
            'AI adoption is moving from experimentation into workflows that affect money, safety, and trust.',
            'The next layer is not another model — it is governance infrastructure around decision-making.',
            'Datacendia defines that layer with multi-agent deliberation plus cryptographic evidence.',
        ],
        'market_title': 'Why the Market Is Opening',
        'market_lines': [
            'Model capability is improving faster than enterprise confidence in unaudited AI decisions.',
            'Buyers increasingly want systems that can be reviewed, challenged, and replayed.',
            'Decision evidence turns AI governance from policy theater into operational software.',
            'The category spans regulated industries, sovereign deployments, and mission-critical decisions.',
            'This is a platform wedge, not a point solution, because the evidence layer sits across workflows.',
            'The more decisions flow through the system, the stronger the switching costs and governance moat.',
            'Open-core distribution can widen the funnel while enterprise controls drive monetization.',
            'The market leader will own the record of why AI-informed decisions were made.',
        ],
        'compliance_title': 'Category Formation',
        'compliance_body': 'Why this matters for fundraising: category-defining companies create a new control plane. Decision evidence can become the default governance layer between enterprise AI deployment and the stakeholders who need to trust it.',
    },
    '10-open-source-strategy.pptx': {
        'cover_label': 'OPEN SOURCE STRATEGY',
        'headline': 'Open core distribution with enterprise-grade governance controls',
        'subheadline': 'Open source can accelerate trust, adoption, and community contribution, while paid controls monetize sovereign deployment, compliance, and enterprise workflow needs.',
        'problem_title': 'Open Source Alone Is Not a Governance Strategy.',
        'problem_lines': [
            'Teams increasingly expect transparent foundations for infrastructure they bet their workflows on.',
            'But open source without evidence, policy, and enterprise controls does not solve accountability.',
            'The opportunity is open distribution plus premium governance capabilities around real deployments.',
            'Few platforms pair community reach with sovereign, audit-ready decision infrastructure.',
            'Developers want inspectable systems; enterprises want control, replay, and support.',
            'That combination can widen top-of-funnel while protecting monetizable enterprise value.',
            'Datacendia can use open core to build trust without giving away the governance moat.',
        ],
        'market_title': 'Why the Open-Core Motion Works',
        'market_lines': [
            'Open foundations can accelerate adoption in developer, research, and public-interest communities.',
            'Enterprise buyers still pay for deployment controls, integrations, support, and compliance workflows.',
            'Governance products benefit from transparency because trust compounds when systems can be inspected.',
            'Decision evidence remains differentiated through workflow depth, deployment breadth, and enterprise operations.',
            'Community usage creates feedback loops, ecosystem integrations, and category awareness.',
            'Commercial layers monetize the teams that need sovereign deployment and defensible controls.',
            'This strategy aligns product trust with distribution leverage.',
            'The result is a larger funnel without collapsing the enterprise moat.',
        ],
        'compliance_title': 'Open-Core Positioning',
        'compliance_body': 'Why this matters for fundraising: open source can strengthen trust and distribution, while the commercial platform captures value from enterprises that need sovereign deployment, workflow governance, and accountable AI operations.',
    },
}


def set_paragraph_text(paragraph, new_text):
    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ''
    else:
        paragraph.text = new_text



def replace_paragraphs(prs, replacements):
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                original = paragraph.text.strip()
                if not original:
                    continue
                for needle, new_text in replacements:
                    if needle in original and original != new_text:
                        set_paragraph_text(paragraph, new_text)
                        changed += 1
                        break
    return changed



def slug_from_filename(filename):
    stem = os.path.splitext(filename)[0]
    return stem.split('-', 1)[1] if '-' in stem else stem



def clean_title(name):
    return name.replace(' & Capital Markets', '').strip()



def design_partner_group(target):
    if target['num'] <= 40:
        return 'vertical'
    if target['num'] <= 55:
        return 'workflow'
    return 'integration'



def design_partner_pack(target):
    name = target['name']
    group = design_partner_group(target)
    if group == 'vertical':
        return {
            'subtitle': f'Auditable AI workflows for {name}',
            'quote': f'"{name} teams are bringing AI into high-stakes workflows. They need evidence, dissent, and accountable approvals — not unexplained outputs."',
            'helps': [
                f'{name} workflow oversight and escalation review',
                'Risk, exception, and approval tracking across teams',
                'Exportable evidence for internal, customer, and regulator review',
            ],
            'scenarios': [
                f'{name} policy and exception review with multi-agent deliberation',
                'Vendor, partner, or procurement evaluation with full evidence trail',
                'Incident replay and root-cause analysis with preserved dissent',
                'Budget and priority setting with accountable approvals',
                'Stakeholder reporting from one searchable decision record',
            ],
        }
    if group == 'workflow':
        return {
            'subtitle': f'Collaborative pilot for {name} with replayable decision evidence',
            'quote': f'"{name} reveals whether an AI workflow is governed or merely automated. High-stakes teams need challenge, evidence, and replay before rollout."',
            'helps': [
                f'{name} workflows with full decision evidence',
                'Human-in-the-loop approvals and preserved dissent',
                'Exportable decision packets for audit, security, and leadership',
            ],
            'scenarios': [
                f'{name} pilot using live scenarios and approval chains',
                'Exception handling and escalation design with traceable outcomes',
                'Decision packet exports for board, regulator, or customer review',
                'Role-based review with replay for every material decision',
                'Roadmap feedback folded directly into product development',
            ],
        }
    return {
        'subtitle': f'Governed AI deployment for {name}',
        'quote': f'"{name} succeeds only when teams can prove what the model saw, who approved the action, and how decisions changed over time."',
        'helps': [
            f'{name} integration with full traceability and replay',
            'Policy controls across cloud, on-prem, and sovereign environments',
            'Evidence exports for security, compliance, and platform teams',
        ],
        'scenarios': [
            f'{name} integration pilot with governed prompts, policies, and approvals',
            'Connector-level access controls and immutable decision history',
            'Change review for models, tools, and workflows before release',
            'Incident replay and rollback with preserved reasoning and dissent',
            'Operations, security, and compliance teams aligned on one evidence chain',
        ],
    }



def enterprise_domain(name):
    lower = name.lower()
    if any(token in lower for token in ['bank', 'credit', 'treasury', 'insurance', 'investment', 'loan', 'trading', 'aml', 'kyc', 'securities', 'wealth', 'pension', 'portfolio', 'tax']):
        return 'finance'
    if any(token in lower for token in ['clinical', 'health', 'biotech', 'drug', 'medical', 'fda', 'hipaa', 'pharma', 'veterinary', 'public health']):
        return 'health'
    if any(token in lower for token in ['defense', 'military', 'election', 'municipal', 'police', 'school', 'university', 'immigration', 'central bank']):
        return 'public'
    if any(token in lower for token in ['energy', 'grid', 'utility', 'construction', 'fleet', 'manufacturing', 'maritime', 'mining', 'oil', 'pipeline', 'transportation', 'supply chain', 'rail', 'water', 'smart building', 'logistics']):
        return 'industrial'
    if any(token in lower for token in ['cyber', 'identity', 'iot', 'zero trust', 'privacy', 'data', 'export control', 'telecom', 'social media', 'vendor risk', 'third-party risk']):
        return 'security'
    return 'general'



def enterprise_use_cases(name):
    domain = enterprise_domain(name)
    if domain == 'finance':
        return [
            'Committee review with preserved dissent and approval history',
            'Exception handling for high-risk workflows with replayable evidence',
            'Policy, risk, and compliance mapping from one evidence chain',
            'Model and workflow change review before production rollout',
            'Board, audit, and regulator reporting with exportable decision packets',
        ]
    if domain == 'health':
        return [
            'Clinical or operational review with preserved dissent and approvals',
            'Exception handling for patient, quality, or compliance workflows',
            'Cross-functional evidence shared across operations, legal, and oversight teams',
            'Model and workflow change review before frontline deployment',
            'Incident replay and stakeholder reporting from one decision record',
        ]
    if domain == 'public':
        return [
            'Policy and program review with structured multi-agent deliberation',
            'High-scrutiny exceptions escalated with accountable approvals',
            'Cross-functional evidence for leadership, legal, and oversight bodies',
            'Vendor, model, and workflow review before operational rollout',
            'Committee, board, or ministerial reporting with exportable packets',
        ]
    if domain == 'industrial':
        return [
            'Operational review with preserved dissent and escalation history',
            'Exception handling for safety, uptime, and service-critical workflows',
            'Cross-team evidence for operations, risk, and compliance stakeholders',
            'Change review for models, automation rules, and frontline workflows',
            'Incident replay and operator reporting from one evidence chain',
        ]
    if domain == 'security':
        return [
            'Security and risk review with preserved reasoning and approvals',
            'Exception handling for controls, alerts, and policy breaches',
            'Evidence shared across engineering, legal, and audit stakeholders',
            'Model and workflow governance before production release',
            'Board, regulator, and customer reporting with exportable packets',
        ]
    return [
        'Decision review with preserved dissent and approval history',
        'Exception handling with role-based escalation and replay',
        'Policy, risk, and compliance mapping from one evidence chain',
        'Vendor, model, or workflow change review before rollout',
        'Board, audit, or operator reporting with exportable decision packets',
    ]



def enterprise_focus(name):
    domain = enterprise_domain(name)
    if domain == 'finance':
        return 'High-stakes financial workflows need replayable reasoning before losses, findings, or regulator questions arrive.'
    if domain == 'health':
        return 'Healthcare and life-science workflows require accountable AI before issues reach patients, quality teams, or compliance review.'
    if domain == 'public':
        return 'Public-interest decisions need challenge, evidence, and human accountability before they face oversight or public scrutiny.'
    if domain == 'industrial':
        return 'Operational AI only scales when teams can reconstruct why an action was taken, who approved it, and how risk was handled.'
    if domain == 'security':
        return 'Security-sensitive AI needs preserved reasoning, accountable approvals, and fast replay when incidents or audits occur.'
    return 'High-stakes AI adoption moves faster when teams can explain, challenge, and replay material decisions.'



def regional_group(name):
    if any(token in name for token in MULTILATERAL_TOKENS):
        return 'multilateral'
    if any(token in name for token in NORTH_AMERICA_TOKENS):
        return 'north-america'
    if any(token in name for token in EUROPE_TOKENS):
        return 'europe'
    if any(token in name for token in ASIA_TOKENS):
        return 'asia'
    if any(token in name for token in AFRICA_TOKENS):
        return 'africa'
    if any(token in name for token in LATAM_TOKENS):
        return 'latam'
    if any(token in name for token in MIDDLE_EAST_TOKENS):
        return 'middle-east'
    return 'global'



def regional_pack(name):
    group = regional_group(name)
    subtitle = f'Decision evidence infrastructure for organizations operating in {name}'
    if group == 'multilateral':
        return {
            'subtitle': subtitle,
            'why_now': [
                'COORDINATION',
                'Multi-jurisdiction programs need a shared record of why decisions were made',
                'GOVERNANCE',
                'Member states and partners need consistent oversight, translation, and replay',
                'EVIDENCE PORTABILITY',
                'Decision evidence must travel across committees, agencies, and funding bodies',
                'Platforms built for cross-border decision evidence',
            ],
            'problem_title': f'Cross-border programs in {name} need accountable AI decisions',
            'problem_body': [
                'Complex institutions often rely on fragmented tools with weak auditability.',
                'When decisions move across agencies, committees, or member states, rationale gets lost.',
                'Leaders need preserved dissent, accountable approvals, and a shared evidence chain.',
                'Without replayable decisions, every review cycle becomes slower, more political, and harder to defend.',
                'AI adoption needs governance infrastructure that works across jurisdictions, languages, and stakeholders.',
            ],
            'solution_intro': f'Datacendia gives teams in {name} multi-agent deliberation, full evidence trails, and deployment flexibility for decisions that must stand up across jurisdictions.',
            'solution_points': [
                'Policy and program review with preserved dissent',
                'Funding, procurement, and partner approvals with exportable evidence',
                'Cross-agency coordination from one shared decision record',
                'Crisis response and after-action replay across stakeholders',
                'Steering committee and board packets in hours, not weeks',
            ],
            'opportunity_title': f'Why {name} Matters',
            'opportunity_body': 'Regional and multilateral coordination increasingly depends on shared evidence, common governance standards, and deployment flexibility across jurisdictions.',
        }
    if group == 'north-america':
        why_now = [
            'AI ADOPTION',
            'Public and private institutions are moving AI from pilots into live operations',
            'OVERSIGHT',
            'Procurement, model risk, and public-trust expectations are rising',
            'EXECUTION',
            'Leaders need replayable decisions before incidents become investigations',
            'Platforms combining sovereign deployment with decision evidence',
        ]
    elif group == 'europe':
        why_now = [
            'AI ADOPTION',
            'European institutions are moving AI into regulated and operational workflows',
            'SOVEREIGNTY',
            'Data residency, procurement, and oversight expectations keep rising',
            'COORDINATION',
            'Cross-border teams need one evidence chain that survives audits and handoffs',
            'Platforms combining sovereignty with decision evidence',
        ]
    elif group == 'asia':
        why_now = [
            'AI ADOPTION',
            'Export-driven and regulated markets are moving AI into production workflows',
            'CONTROL',
            'Data residency, supplier risk, and multilingual operations raise governance complexity',
            'EXECUTION',
            'Decision evidence helps teams move faster without losing accountability',
            'Platforms built for accountable AI deployment',
        ]
    elif group == 'africa':
        why_now = [
            'AI ADOPTION',
            'Public institutions and growth-stage enterprises are moving AI into frontline workflows',
            'ACCOUNTABILITY',
            'Cross-border coordination and stakeholder scrutiny are increasing',
            'EXECUTION',
            'Leaders need accountable decisions without heavyweight data movement',
            'Platforms built for decision evidence and flexible deployment',
        ]
    elif group == 'latam':
        why_now = [
            'AI ADOPTION',
            'Regional enterprises and institutions are modernizing decisions with AI',
            'ACCOUNTABILITY',
            'Cross-border business and regulatory expectations create governance pressure',
            'EXECUTION',
            'Teams need evidence-rich workflows before trust breaks under scrutiny',
            'Platforms built for accountable AI across complex operating environments',
        ]
    elif group == 'middle-east':
        why_now = [
            'AI ADOPTION',
            'National transformation programs are accelerating AI deployment across sectors',
            'SOVEREIGNTY',
            'Critical infrastructure and sensitive workflows demand strong deployment control',
            'COORDINATION',
            'Cross-ministry and enterprise teams need one defensible evidence chain',
            'Platforms combining sovereign deployment with decision evidence',
        ]
    else:
        why_now = [
            'AI ADOPTION',
            'Regional institutions are moving AI into regulated and operational workflows',
            'OVERSIGHT',
            'Risk, legal, and procurement expectations are increasing',
            'EXECUTION',
            'Teams need replayable decisions before issues become investigations',
            'Platforms built for accountable AI deployment',
        ]
    return {
        'subtitle': subtitle,
        'why_now': why_now,
        'problem_title': f'Organizations across {name} are deploying AI without decision evidence',
        'problem_body': [
            'AI is moving into high-stakes workflows faster than governance is catching up.',
            'Public and private institutions need to explain who approved decisions, what evidence was used, and how dissent was handled.',
            'Fragmented tooling makes it hard to replay decisions across teams, committees, and operating environments.',
            'Without accountable approvals and preserved reasoning, trust erodes when a decision is challenged.',
            'Decision evidence is the missing layer between AI deployment and defensible regional execution.',
        ],
        'solution_intro': f'Datacendia gives teams in {name} multi-agent deliberation, full evidence trails, and sovereign deployment for decisions that must stand up to scrutiny.',
        'solution_points': [
            'Regulatory and policy review with preserved dissent',
            'Vendor, model, and procurement approvals with exportable evidence',
            'Board, ministry, or committee briefing from one evidence chain',
            'Cross-border coordination with accountable approvals and replay',
            'Incident response and after-action review in minutes, not weeks',
        ],
        'opportunity_title': f'Why {name} Matters',
        'opportunity_body': 'Regional buyers need AI systems that respect local oversight, data-control expectations, and real-world operating complexity — not generic black-box automation.',
    }



def fix_design_partner_decks():
    changed_files = 0
    for filename in sorted(os.listdir(os.path.join(BASE, 'design-partner', 'pptx'))):
        if not filename.endswith('.pptx'):
            continue
        slug = slug_from_filename(filename)
        target = DESIGN_PARTNER_BY_SLUG.get(slug)
        if not target:
            continue
        pack = design_partner_pack(target)
        path = os.path.join(BASE, 'design-partner', 'pptx', filename)
        prs = Presentation(path)
        replacements = [
            ('Clinical decision support with HIPAA-compliant audit trails', pack['subtitle']),
            ('Clinical decisions need explainable, auditable AI', pack['quote']),
            ('patient safety as the prime directive', pack['quote']),
            ('Clinical systems integration and health informatics', pack['helps'][0]),
            ('Safety metrics, adverse event prevention', pack['helps'][1]),
            ('HIPAA, FDA, CMS regulatory compliance', pack['helps'][2]),
            ('Clinical protocol review — Multi-agent evaluation against safety, regulatory, and financial criteria', pack['scenarios'][0]),
            ('Hospital M&A due diligence — Clinical outcomes, compliance, financial health analysis', pack['scenarios'][1]),
            ('Regulatory change impact — CMS/FDA guidance mapped across operations', pack['scenarios'][2]),
            ('Patient safety incident review — Full decision replay with evidence preservation', pack['scenarios'][3]),
            ('Value-based care optimization — Balance quality, outcomes, and financial sustainability', pack['scenarios'][4]),
        ]
        if replace_paragraphs(prs, replacements):
            prs.save(path)
            changed_files += 1
    return changed_files



def fix_enterprise_decks():
    changed_files = 0
    for filename in sorted(os.listdir(os.path.join(BASE, 'enterprise', 'pptx'))):
        if not filename.endswith('.pptx'):
            continue
        slug = slug_from_filename(filename)
        target = ENTERPRISE_BY_SLUG.get(slug)
        if not target:
            continue
        name = target['name']
        focus = enterprise_focus(name)
        intro = f'Multi-agent deliberation with cryptographic evidence for every {name.lower()} decision — review-ready before the committee, customer, or regulator asks.'
        use_cases = enterprise_use_cases(name)
        path = os.path.join(BASE, 'enterprise', 'pptx', filename)
        prs = Presentation(path)
        replacements = [
            (' & Capital Markets', ''),
            ('AI Governance for Regulated Financial Institutions', f'Accountable AI workflows for {name}'),
            ('DORA', 'AI ADOPTION'),
            ('Jan 2025 — AI Operational Resilience', f'{name} workflows are moving from pilot to production'),
            ('Basel III', 'ACCOUNTABILITY'),
            ('Auditable Risk Models Required', 'Leaders need replayable evidence for every material decision'),
            ('MiFID II', 'EXECUTION'),
            ('Full Audit Trail for Algo Trading', 'Operations need challenge, approval, and audit — not black boxes'),
            ('AI Platforms w/ Crypto Decision Evidence', 'Platforms built around decision evidence by design'),
            ('Banks Deploy AI Without Audit Trails', f'Organizations use AI in {name} without decision evidence'),
            ('Credit decisions influenced by AI with no explainability for regulators', 'Teams automate recommendations, approvals, or exception handling with limited traceability'),
            ('DORA requires operational resilience documentation for all AI systems by Jan 2025', 'When outcomes are challenged, the rationale behind the model and human overrides is hard to reconstruct'),
            ('Basel III mandates auditable risk models — AI black boxes fail this test', 'Legal, risk, and leadership teams need preserved dissent, accountable approvals, and replay'),
            ('MiFID II requires full audit trails for algorithmic trading decisions', 'Siloed tooling makes it expensive to explain who decided what, why, and with what evidence'),
            ('No widely adopted platform provides cryptographically verifiable decision evidence at the decision level', 'The result: slower adoption, higher exposure, and no clean audit trail for high-stakes decisions'),
            ('Multi-agent deliberation with cryptographic evidence for every credit, trading, and risk decision — regulator-ready proof before the auditor calls.', intro),
            ('Credit decision review with dissent tracking', use_cases[0]),
            ('M&A due diligence with full evidence trail', use_cases[1]),
            ('DORA compliance with automated evidence', use_cases[2]),
            ('Investment committee structured deliberation', use_cases[3]),
            ('Regulatory change mapping to P&L impact', use_cases[4]),
            ('INVESTOR THESIS', 'WHY THIS USE CASE'),
            ('Why This Is a Compelling Investment', f'Why {name} Matters'),
            ('Only platform providing cryptographic decision evidence — a regulatory requirement with zero incumbent solutions.', f'AI adoption in {name} is accelerating, but trust breaks when teams cannot explain outcomes. Datacendia adds multi-agent challenge, evidence, and replay before the audit, incident review, or board meeting. {focus}'),
            ('Built the exact data infrastructure that financial institutions need for AI governance — at one of the world\'s largest asset managers.', 'Built the data infrastructure and decision workflows required for regulated, high-stakes AI deployments at institutional scale.'),
        ]
        if replace_paragraphs(prs, replacements):
            prs.save(path)
            changed_files += 1
    return changed_files



def fix_regional_decks():
    changed_files = 0
    for filename in sorted(os.listdir(os.path.join(BASE, 'regional', 'pptx'))):
        if not filename.endswith('.pptx'):
            continue
        slug = slug_from_filename(filename)
        target = REGIONAL_BY_SLUG.get(slug)
        if not target:
            continue
        pack = regional_pack(target['name'])
        path = os.path.join(BASE, 'regional', 'pptx', filename)
        prs = Presentation(path)
        replacements = [
            (' & Capital Markets', ''),
            ('AI Governance for Regulated Financial Institutions', pack['subtitle']),
            ('DORA', pack['why_now'][0]),
            ('Jan 2025 — AI Operational Resilience', pack['why_now'][1]),
            ('Basel III', pack['why_now'][2]),
            ('Auditable Risk Models Required', pack['why_now'][3]),
            ('MiFID II', pack['why_now'][4]),
            ('Full Audit Trail for Algo Trading', pack['why_now'][5]),
            ('AI Platforms w/ Crypto Decision Evidence', pack['why_now'][6]),
            ('Banks Deploy AI Without Audit Trails', pack['problem_title']),
            ('Credit decisions influenced by AI with no explainability for regulators', pack['problem_body'][0]),
            ('DORA requires operational resilience documentation for all AI systems by Jan 2025', pack['problem_body'][1]),
            ('Basel III mandates auditable risk models — AI black boxes fail this test', pack['problem_body'][2]),
            ('MiFID II requires full audit trails for algorithmic trading decisions', pack['problem_body'][3]),
            ('No widely adopted platform provides cryptographically verifiable decision evidence at the decision level', pack['problem_body'][4]),
            ('Multi-agent deliberation with cryptographic evidence for every credit, trading, and risk decision — regulator-ready proof before the auditor calls.', pack['solution_intro']),
            ('Credit decision review with dissent tracking', pack['solution_points'][0]),
            ('M&A due diligence with full evidence trail', pack['solution_points'][1]),
            ('DORA compliance with automated evidence', pack['solution_points'][2]),
            ('Investment committee structured deliberation', pack['solution_points'][3]),
            ('Regulatory change mapping to P&L impact', pack['solution_points'][4]),
            ('INVESTOR THESIS', 'REGIONAL OPPORTUNITY'),
            ('Why This Is a Compelling Investment', pack['opportunity_title']),
            ('Only platform providing cryptographic decision evidence — a regulatory requirement with zero incumbent solutions.', pack['opportunity_body']),
            ('Built the exact data infrastructure that financial institutions need for AI governance — at one of the world\'s largest asset managers.', 'Built the data infrastructure and decision workflows required for regulated, high-stakes AI deployments at institutional scale.'),
        ]
        if replace_paragraphs(prs, replacements):
            prs.save(path)
            changed_files += 1
    return changed_files



def fix_gamma_decks():
    changed_files = 0
    gamma_dir = os.path.join(BASE, 'gamma', 'pptx')
    overview_path = os.path.join(gamma_dir, '01-investor-overview.pptx')
    overview = Presentation(overview_path)
    if replace_paragraphs(overview, [('Decision CrisisImmunizationInfrastructure', 'Decision Evidence Infrastructure for AI Systems')]):
        overview.save(overview_path)
        changed_files += 1

    for filename, pack in GAMMA_CONTENT.items():
        path = os.path.join(gamma_dir, filename)
        prs = Presentation(path)
        replacements = [
            ('SERIES A', pack['cover_label']),
            ('Decision CrisisImmunizationInfrastructure', pack['headline']),
            ('50+ specialized AI agent presets debate, dissent, and decide — with cryptographic proof of every step. Not a chatbot. A boardroom.', pack['subheadline']),
            ('AI Makes Decisions.Nobody Can Explain Them.', pack['problem_title']),
            ('$7.9T in decisions made using AI by 2030 with no audit trail', pack['problem_lines'][0]),
            ('73% of enterprises cite AI explainability as top concern (Gartner)', pack['problem_lines'][1]),
            ('$2.1B in regulatory fines for unexplainable algorithmic decisions (2023)', pack['problem_lines'][2]),
            ('Zero platforms provide multi-agent deliberation with cryptographic proof', pack['problem_lines'][3]),
            ('Enterprises are deploying AI at scale — but every decision is a black box. When the regulator asks "why did your AI do that?", the answer is silence.', pack['problem_lines'][4]),
            ('Single-model AI gives you one answer. Real decisions require debate, dissent, and documented reasoning. Boards don\'t accept one person\'s opinion — why accept one model\'s?', pack['problem_lines'][5]),
            ('Datacendia fills the gap between "AI generated an answer" and "AI made a defensible decision."', pack['problem_lines'][6]),
            ('Category-Creating in a $50B+ Market', pack['market_title']),
            ('EU AI Act — mandates explainability for high-risk AI systems (2026 enforcement)', pack['market_lines'][0]),
            ('US Executive Order on AI — federal agencies must document AI decision processes', pack['market_lines'][1]),
            ('DORA (Digital Operational Resilience Act) — financial services AI governance mandate', pack['market_lines'][2]),
            ('Enterprise AI spend — projected $200B+ by 2028, governance is lagging', pack['market_lines'][3]),
            ('Regulation is arriving — first-mover advantage in compliance tooling', pack['market_lines'][4]),
            ('LLMs made multi-agent possible — viable at enterprise scale since 2024', pack['market_lines'][5]),
            ('Enterprise AI budgets unlocked — every F500 has an AI governance gap', pack['market_lines'][6]),
            ('Trust deficit growing — AI hallucinations eroding confidence', pack['market_lines'][7]),
            ('Enterprise Compliance Roadmap', pack['compliance_title']),
            ('Why this matters for fundraising: Every compliance framework unlocks a market segment. SOC 2 → enterprise SaaS. HIPAA → healthcare ($4T market). FedRAMP → federal government ($100B+ IT spend). Each certification is a sales accelerant, not just a checkbox.', pack['compliance_body']),
        ]
        if replace_paragraphs(prs, replacements):
            prs.save(path)
            changed_files += 1
    return changed_files



def fix_investor_50_index_screenshot():
    path = os.path.join(BASE, 'investor-50', 'pptx', 'index.pptx')
    prs = Presentation(path)
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
    if image_count > 0:
        return 0
    screenshot = os.path.join(SCREENSHOTS, '02-council-deliberation.png')
    slide = prs.slides[0]
    slide.shapes.add_picture(screenshot, Inches(7.0), Inches(1.2), width=Inches(5.3))
    prs.save(path)
    return 1



def main():
    print('Resolving review issues...')
    design_partner = fix_design_partner_decks()
    enterprise = fix_enterprise_decks()
    regional = fix_regional_decks()
    gamma = fix_gamma_decks()
    index_fixed = fix_investor_50_index_screenshot()
    print(f'  design-partner fixed: {design_partner}')
    print(f'  enterprise fixed:     {enterprise}')
    print(f'  regional fixed:       {regional}')
    print(f'  gamma fixed:          {gamma}')
    print(f'  index screenshots:    {index_fixed}')
    print('Done.')


if __name__ == '__main__':
    main()
