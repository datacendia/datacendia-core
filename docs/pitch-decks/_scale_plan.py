"""
Plan for scaling from 240 to 1,000 pitch decks.
First: extract the EXACT structure of an existing investor-100 deck as a template.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Emu, Pt
import os, json

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

# Use a16z as the reference template
template = os.path.join(BASE, 'investor-100', 'pptx', '001-a16z.pptx')
prs = Presentation(template)

print(f'Slide dimensions: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inches')
print(f'Total slides: {len(prs.slides)}\n')

for i, slide in enumerate(prs.slides):
    print(f'{"="*60}')
    print(f'SLIDE {i+1}')
    print(f'{"="*60}')
    
    for shape in slide.shapes:
        stype = shape.shape_type
        print(f'\n  Shape: {shape.shape_type} at ({shape.left/914400:.2f},{shape.top/914400:.2f}) size {shape.width/914400:.2f}x{shape.height/914400:.2f}')
        
        if stype == MSO_SHAPE_TYPE.PICTURE:
            print(f'    [IMAGE]')
        
        if shape.has_text_frame:
            for j, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    # Get font info
                    font_info = ''
                    if para.runs:
                        run = para.runs[0]
                        size = run.font.size
                        bold = run.font.bold
                        color = run.font.color.rgb if run.font.color and run.font.color.rgb else None
                        font_info = f' [size={size}, bold={bold}, color={color}]'
                    print(f'    P{j}{font_info}: "{t[:100]}"')

# Current counts
print('\n\n' + '='*60)
print('CURRENT DECK INVENTORY')
print('='*60)
current = {
    'main': 3,
    'investor': 10,
    'investor-50': 51,
    'investor-100': 120,
    'sales': 25,
    'design-partner': 25,
    'gamma': 6,
}
print(f'Current total: {sum(current.values())}')

# Target: 1,000
target = {
    'main': 3,           # +0 (no change)
    'investor': 15,       # +5 (more generic themes)
    'investor-50': 100,   # +49 (more thesis/vertical decks)
    'investor-500': 500,  # +380 (rename investor-100 -> investor-500, add 380 real firms)
    'sales': 75,          # +50 (more verticals, pain points, deployment scenarios)
    'design-partner': 75, # +50 (more industries, use cases)
    'gamma': 10,          # +4 (more presentation formats)
    'enterprise': 100,    # NEW: industry-specific enterprise target decks
    'regional': 122,      # NEW: region/country-specific decks
}
print(f'Target total: {sum(target.values())}')
print(f'New decks needed: {sum(target.values()) - sum(current.values())}')

for cat, count in target.items():
    old = current.get(cat, 0)
    diff = count - old
    if diff > 0:
        print(f'  {cat}: {old} -> {count} (+{diff})')
