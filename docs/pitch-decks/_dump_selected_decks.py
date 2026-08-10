import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'c:\Users\Stu\datacendia-core\docs\pitch-decks'
FILES = [
    'design-partner/pptx/26-real-estate.pptx',
    'design-partner/pptx/41-pre-mortem.pptx',
    'enterprise/pptx/50-logistics-optimization.pptx',
    'regional/pptx/39-lithuania-vilnius.pptx',
    'gamma/pptx/07-healthcare-vertical.pptx',
    'gamma/pptx/08-financial-compliance.pptx',
]
for rel in FILES:
    path = os.path.join(BASE, rel)
    print('\n' + '='*90)
    print(rel)
    print('='*90)
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        print(f'\n--- Slide {i} ---')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        print(' ', t)
