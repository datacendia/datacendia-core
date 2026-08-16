"""Extract all text from templates - UTF-8 output."""
import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

BASE = os.path.dirname(__file__)
path = os.path.join(BASE, 'investor-100', 'pptx', '001-a16z.pptx')
prs = Presentation(path)

for si, slide in enumerate(prs.slides):
    print(f'\n--- Slide {si+1} ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  {text}')
        if shape.shape_type == 13:
            print(f'  [IMAGE] at ({shape.left},{shape.top})')
