"""Comprehensive review of all pitch decks for text quality and category health."""
import os, sys, io, glob, hashlib
from collections import defaultdict, Counter
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = os.path.dirname(os.path.abspath(__file__))
FOLDERS = {
    'pptx': 'main',
    'investor/pptx': 'investor',
    'investor-50/pptx': 'investor-50',
    'investor-100/pptx': 'investor-100',
    'investor-500/pptx': 'investor-500',
    'sales/pptx': 'sales',
    'design-partner/pptx': 'design-partner',
    'gamma/pptx': 'gamma',
    'enterprise/pptx': 'enterprise',
    'regional/pptx': 'regional',
}

BANNED_SNIPPETS = [
    '90-day guarantee',
    'warm enterprise intros',
    'platforms w/ crypto ai evidence',
    '0 platforms w/ crypto ai evidence',
    'ai governance infrastructure',
    'missing governance layer for enterprise ai',
]

REQUIRED_SNIPPETS = [
    'decision evidence',
    'active enterprise design-partner conversations',
]

CATEGORY_EXPECTED_SLIDES = {
    'main': {10},
    'investor': {10},
    'investor-50': {10},
    'investor-100': {10},
    'investor-500': {10},
    'sales': {8},
    'design-partner': {9},
    'gamma': {10},
    'enterprise': {10},
    'regional': {10},
}


def deck_text(prs):
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
    return '\n'.join(parts)


def find_images(prs):
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
    return count


def first_meaningful_texts(prs, max_items=8):
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and t != 'DATACENDIA' and 'datacendia.com' not in t.lower() and '© 2024' not in t:
                        out.append(t)
                        if len(out) >= max_items:
                            return out
    return out


summary = []
banned_hits = []
missing_required = []
slide_anomalies = []
no_images = []
hash_map = defaultdict(list)
category_counts = Counter()
category_slide_counts = defaultdict(Counter)
samples = {}

for folder, category in FOLDERS.items():
    path = os.path.join(BASE, folder)
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    category_counts[category] = len(files)
    if not files:
        continue

    # sample a few representative files per category
    sample_candidates = [files[0], files[len(files)//2], files[-1]] if len(files) >= 3 else files
    unique_samples = []
    seen = set()
    for f in sample_candidates:
        if f not in seen:
            unique_samples.append(f)
            seen.add(f)
    samples[category] = unique_samples

    for f in files:
        rel = os.path.relpath(f, BASE)
        try:
            prs = Presentation(f)
        except Exception as e:
            summary.append((rel, f'ERROR opening: {e}'))
            continue

        text = deck_text(prs).lower()
        image_count = find_images(prs)
        slide_count = len(prs.slides)
        category_slide_counts[category][slide_count] += 1

        if slide_count not in CATEGORY_EXPECTED_SLIDES[category]:
            slide_anomalies.append((rel, slide_count, category))

        if image_count == 0:
            no_images.append(rel)

        for banned in BANNED_SNIPPETS:
            if banned in text:
                banned_hits.append((rel, banned))

        for required in REQUIRED_SNIPPETS:
            if required not in text and category in {'investor-100', 'investor-500'}:
                missing_required.append((rel, required))

        h = hashlib.md5(open(f, 'rb').read()).hexdigest()
        hash_map[h].append(rel)

exact_duplicates = [paths for paths in hash_map.values() if len(paths) > 1]

print('='*80)
print('CATEGORY COUNTS')
print('='*80)
for cat, count in sorted(category_counts.items()):
    print(f'{cat:15} {count}')

print('\n' + '='*80)
print('SLIDE COUNT DISTRIBUTION BY CATEGORY')
print('='*80)
for cat in sorted(category_slide_counts):
    print(f'\n{cat}:')
    for slide_count, freq in sorted(category_slide_counts[cat].items()):
        print(f'  {slide_count} slides -> {freq} deck(s)')

print('\n' + '='*80)
print('BANNED COPY HITS')
print('='*80)
if banned_hits:
    for rel, banned in banned_hits[:100]:
        print(f'  {rel}: {banned}')
else:
    print('  None found')

print('\n' + '='*80)
print('MISSING REQUIRED INVESTOR COPY')
print('='*80)
if missing_required:
    for rel, req in missing_required[:100]:
        print(f'  {rel}: missing "{req}"')
else:
    print('  None found')

print('\n' + '='*80)
print('SLIDE COUNT ANOMALIES')
print('='*80)
if slide_anomalies:
    for rel, slides, cat in slide_anomalies[:100]:
        print(f'  {rel}: {slides} slides (category {cat})')
else:
    print('  None found')

print('\n' + '='*80)
print('DECKS WITH NO IMAGES')
print('='*80)
if no_images:
    for rel in no_images[:100]:
        print(f'  {rel}')
else:
    print('  None found')

print('\n' + '='*80)
print('EXACT DUPLICATE FILES')
print('='*80)
filtered_dupes = []
for paths in exact_duplicates:
    if len(paths) > 1:
        filtered_dupes.append(paths)
if filtered_dupes:
    for paths in filtered_dupes[:20]:
        print('  DUPLICATE SET:')
        for p in paths:
            print(f'    {p}')
else:
    print('  None found')

print('\n' + '='*80)
print('REPRESENTATIVE SAMPLES')
print('='*80)
for cat, files in samples.items():
    print(f'\n[{cat}]')
    for f in files:
        rel = os.path.relpath(f, BASE)
        try:
            prs = Presentation(f)
            texts = first_meaningful_texts(prs)
            print(f'  {rel}')
            for t in texts[:6]:
                print(f'    - {t}')
        except Exception as e:
            print(f'  {rel} ERROR: {e}')
