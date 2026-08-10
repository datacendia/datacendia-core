"""Check remaining template bleed issues across investor-500 decks."""
import os, sys, glob, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

files = sorted(glob.glob(os.path.join(BASE, 'investor-500', 'pptx', '*.pptx')))

checks = {
    'dedicated AI fund': 0,
    'AI is eating software': 0,
    'has backed the data stack': 0,
    'Anyscale': 0,
    'dbt Labs': 0,
    'Replit': 0,
}

# Also check for slug-as-name issues
slug_issues = 0

for f in files:
    try:
        prs = Presentation(f)
    except:
        continue
    text = '\n'.join(
        p.text for sl in prs.slides for sh in sl.shapes
        if sh.has_text_frame for p in sh.text_frame.paragraphs
    )
    for key in checks:
        if key in text:
            checks[key] += 1

    # Check for slug used as name (hyphens in display text)
    fname = os.path.basename(f)
    slug = fname.replace('.pptx', '').split('-', 1)[1] if '-' in fname else ''
    if slug and len(slug) > 3 and '-' in slug:
        # slugs with hyphens shouldn't appear in running text
        if slug + "'s" in text.lower() or slug + ' has' in text.lower() or slug + ' invests' in text.lower():
            slug_issues += 1

print(f'Total files: {len(files)}')
print()
for key, count in checks.items():
    flag = ' !!!' if count > 0 else ''
    print(f'  "{key}": {count}/{len(files)}{flag}')
print(f'  Slug-as-name: {slug_issues}/{len(files)}')

# Also check Hex separately (it's a short word that could be a false positive)
hex_count = 0
for f in files:
    try:
        prs = Presentation(f)
    except:
        continue
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if t == 'Hex':  # exact match as a bullet point
                    hex_count += 1
                    break
print(f'  "Hex" as bullet: {hex_count}/{len(files)}')
