"""
Deep audit: for every deck, show the FULL text content of slides that have images.
This tells us what screenshot should replace each placeholder.
Also show the exact image position/size so we can match it.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os, glob, json

BASE = r'C:\Users\Stu\datacendia-core\docs\pitch-decks'

folders = ['pptx', 'investor/pptx', 'sales/pptx', 'design-partner/pptx', 
           'gamma/pptx', 'investor-50/pptx', 'investor-100/pptx']

# For each category, examine the image slides of ALL decks
for folder in folders:
    path = os.path.join(BASE, folder)
    if not os.path.exists(path):
        continue
    files = sorted(glob.glob(os.path.join(path, '*.pptx')))
    
    print(f'\n{"="*80}')
    print(f'{folder}/ — {len(files)} decks')
    print(f'{"="*80}')
    
    for f in files:
        try:
            prs = Presentation(f)
            name = os.path.basename(f)
            
            # Find slides with images
            img_slides = []
            for i, slide in enumerate(prs.slides):
                has_img = False
                img_info = None
                texts = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        has_img = True
                        img_info = {
                            'w': round(shape.width / 914400, 2),
                            'h': round(shape.height / 914400, 2),
                            'left': round(shape.left / 914400, 2),
                            'top': round(shape.top / 914400, 2),
                        }
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and t != 'DATACENDIA' and 'datacendia.com' not in t.lower() and '2024' not in t and '2026' not in t:
                                texts.append(t)
                if has_img:
                    img_slides.append({
                        'slide_num': i + 1,
                        'img': img_info,
                        'texts': texts
                    })
            
            if img_slides:
                print(f'\n  {name}')
                for s in img_slides:
                    img = s['img']
                    print(f'    Slide {s["slide_num"]} — img at ({img["left"]},{img["top"]}) size {img["w"]}x{img["h"]}in')
                    for t in s['texts'][:5]:  # first 5 text elements
                        print(f'      "{t[:100]}"')
        except Exception as e:
            print(f'  ERROR {name}: {e}')
