import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation

f = r'C:\Users\Stu\datacendia-core\docs\pitch-decks\investor-500\pptx\378-kuang-chi.pptx'
prs = Presentation(f)
changed = False
for sl in prs.slides:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if 'kuang-chi' in r.text.lower():
                    old = r.text
                    r.text = r.text.replace("kuang-chi's", "Kuang-Chi's")
                    r.text = r.text.replace("kuang-chi has", "Kuang-Chi has")
                    r.text = r.text.replace("kuang-chi invests", "Kuang-Chi invests")
                    r.text = r.text.replace("kuang-chi investment", "Kuang-Chi investment")
                    if r.text != old:
                        changed = True
                        print(f'  Fixed: {old[:80]} -> {r.text[:80]}')
if changed:
    prs.save(f)
    print('Saved.')
else:
    print('No changes needed.')
